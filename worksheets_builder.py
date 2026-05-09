"""
worksheets_builder.py
Importable wrapper around the python-pptx worksheets builder.
Call build_worksheets(lesson_dict) and get back PPTX bytes.
"""

import io
import json
import os
import sys
from pathlib import Path

# Add parent if needed
sys.path.insert(0, str(Path(__file__).parent))


def build_worksheets(lesson: dict) -> bytes:
    """
    Build the 7-slide worksheets PPTX from a lesson dict.
    Returns raw PPTX bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    WORDS       = lesson["words"]
    DEFS        = lesson["defs"]
    SENTENCES   = lesson["sentences"]
    SPELL_DATA  = lesson["spellData"]
    CLOZE_ORDER = lesson["clozeOrder"]
    CODE        = lesson["code"]

    SW = 10.0
    SH = 5.625
    MARGIN    = 0.35
    CONTENT_Y = 0.95
    CONTENT_H = SH - CONTENT_Y - 0.18
    CW        = SW - MARGIN * 2

    FONT = "Aptos"

    C = dict(
        BLACK   = "1A1A1A",
        DGREY   = "444444",
        GREY    = "999999",
        LGREY   = "F2F2F2",
        BGREY   = "CCCCCC",
        PURPLE  = "6750A4",
        LPURPLE = "EDE7F6",
        PINK    = "E91E8C",
        GREEN   = "2E7D32",
        WHITE   = "FFFFFF",
        BDR     = "AAAAAA",
        RED     = "C0392B",
        BOX1    = "4A3F9F",
    )

    def rgb(hex_str):
        return RGBColor.from_string(hex_str)

    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    BLANK = prs.slide_layouts[6]

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _set_bodyPr(txBody, valign="middle", margin_in=0.05):
        bp = txBody.find(qn('a:bodyPr'))
        if bp is None:
            bp = etree.SubElement(txBody, qn('a:bodyPr'))
        bp.set('anchor', {"top": "t", "middle": "ctr", "bottom": "b"}.get(valign, "ctr"))
        m = str(Inches(margin_in))
        for attr in ('lIns', 'rIns', 'tIns', 'bIns'):
            bp.set(attr, m)

    def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_pt=0, radius=None):
        shape = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_hex)
        ln = shape.line
        if line_hex and line_pt > 0:
            ln.color.rgb = rgb(line_hex)
            ln.width = Pt(line_pt)
        else:
            ln.fill.background()
        if radius is not None:
            sp = shape._element
            spPr = sp.find(qn('p:spPr'))
            prstGeom = spPr.find(qn('a:prstGeom'))
            if prstGeom is None:
                prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            for av in list(avLst):
                avLst.remove(av)
            shortest = min(w, h)
            adj_val = int(min(radius / shortest, 0.5) * 100000)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {adj_val}')
        return shape

    def add_text(slide, text, x, y, w, h, *,
                 font=FONT, size=11, bold=False, italic=False,
                 color="1A1A1A", align="left", valign="middle",
                 wrap=True, margin_in=0.05):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        _set_bodyPr(tf._txBody, valign=valign, margin_in=margin_in)
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = text
        rf = run.font
        rf.name = font
        rf.size = Pt(size)
        rf.bold = bold
        rf.italic = italic
        rf.color.rgb = rgb(color)
        return txBox

    def add_rich_text(slide, runs, x, y, w, h, *,
                      font=FONT, align="left", valign="middle", margin_in=0.05):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        _set_bodyPr(tf._txBody, valign=valign, margin_in=margin_in)
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER}.get(align, PP_ALIGN.LEFT)
        for text, opts in runs:
            run = p.add_run()
            run.text = text
            rf = run.font
            rf.name = font
            rf.size = Pt(opts.get("size", 11))
            rf.bold = opts.get("bold", False)
            rf.italic = opts.get("italic", False)
            rf.color.rgb = rgb(opts.get("color", "1A1A1A"))
        return txBox

    def add_wlines(slide, x, y, w, n, spacing, color="6750A4"):
        for i in range(n):
            add_rect(slide, x, y + i * spacing, w, 0.013, color)

    def add_header(slide, activity, instruction, sv_label=""):
        label_right = f"{activity}  |  {sv_label}" if sv_label else activity
        add_text(slide, label_right,
                 MARGIN, 0.10, SW - MARGIN * 2, 0.24,
                 size=8.5, color=C["GREY"], align="left", valign="middle", margin_in=0)
        add_text(slide, instruction,
                 MARGIN, 0.32, SW - MARGIN * 2 - 2.2, 0.50,
                 size=11.5, bold=True, color=C["BLACK"], align="left", valign="middle", margin_in=0)
        add_rect(slide, MARGIN, 0.84, SW - MARGIN * 2, 0.020, C["BGREY"])

    def add_table(slide, rows, x, y, w, h, col_widths, row_heights,
                  border_hex="AAAAAA", border_pt=1.0):
        n_rows = len(rows)
        n_cols = len(rows[0])
        tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                            Inches(x), Inches(y), Inches(w), Inches(h))
        tbl = tbl_shape.table
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)
        rh_list = row_heights if isinstance(row_heights, list) else [row_heights] * n_rows
        for ri, rh in enumerate(rh_list):
            tbl.rows[ri].height = Inches(rh)

        for ri, row in enumerate(rows):
            for ci, cd in enumerate(row):
                cell = tbl.cell(ri, ci)
                text   = cd.get("text", "")
                bold   = cd.get("bold", False)
                italic = cd.get("italic", False)
                color  = cd.get("color", "1A1A1A")
                fill   = cd.get("fill", "FFFFFF")
                align  = cd.get("align", "left")
                size   = cd.get("size", 11)
                valign = cd.get("valign", "middle")

                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(fill)

                tf = cell.text_frame
                tf.word_wrap = True
                _set_bodyPr(tf._txBody, valign=valign, margin_in=0.04)

                p = tf.paragraphs[0]
                p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                               "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
                if text:
                    run = p.add_run()
                    run.text = text
                    rf = run.font
                    rf.name = FONT
                    rf.size = Pt(size)
                    rf.bold = bold
                    rf.italic = italic
                    rf.color.rgb = rgb(color)

                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
                    ln_el = tcPr.find(qn(f'a:{edge}'))
                    if ln_el is None:
                        ln_el = etree.SubElement(tcPr, qn(f'a:{edge}'))
                    ln_el.set('w', str(int(Pt(border_pt))))
                    sf = ln_el.find(qn('a:solidFill'))
                    if sf is None:
                        sf = etree.SubElement(ln_el, qn('a:solidFill'))
                    srgb = sf.find(qn('a:srgbClr'))
                    if srgb is None:
                        srgb = etree.SubElement(sf, qn('a:srgbClr'))
                    srgb.set('val', border_hex)
        return tbl_shape

    # ── Slide builders ────────────────────────────────────────────────────────

    def build_slide1():
        slide = prs.slides.add_slide(BLANK)
        add_header(slide, "Word Sort",
                   "Can you sort this week's words? Some can ONLY be a verb. Others can be a verb OR a noun.",
                   f"{CODE}.6")
        chip_w, chip_h, chip_gap = 0.81, 0.34, 0.042
        n = len(WORDS)
        total_w = n * chip_w + (n - 1) * chip_gap
        chip_x0 = MARGIN + (CW - total_w) / 2
        chip_y  = CONTENT_Y + 0.05
        for i, word in enumerate(WORDS):
            cx = chip_x0 + i * (chip_w + chip_gap)
            add_rect(slide, cx, chip_y, chip_w, chip_h, C["LGREY"],
                     line_hex=C["BGREY"], line_pt=0.8, radius=0.07)
            add_text(slide, word, cx, chip_y, chip_w, chip_h,
                     size=10.5, bold=True, color=C["BLACK"],
                     align="center", valign="middle", margin_in=0)
        box_y = CONTENT_Y + 0.50
        box_h = SH - box_y - 0.15
        box_w = (CW - 0.30) / 2
        ws = lesson["wordSort"]
        boxes = [
            dict(x=MARGIN,              bdr=C["BOX1"], label=ws["box1label"], sub=ws["box1sub"], lines=4),
            dict(x=MARGIN + box_w + 0.30, bdr=C["GREEN"], label=ws["box2label"], sub=ws["box2sub"], lines=6),
        ]
        for b in boxes:
            add_rect(slide, b["x"], box_y, box_w, box_h, C["WHITE"],
                     line_hex=b["bdr"], line_pt=2, radius=0.10)
            add_text(slide, b["label"],
                     b["x"] + 0.12, box_y + 0.08, box_w - 0.24, 0.30,
                     size=12.5, bold=True, color=b["bdr"], align="center", margin_in=0)
            add_text(slide, b["sub"],
                     b["x"] + 0.12, box_y + 0.35, box_w - 0.24, 0.22,
                     size=9.5, italic=True, color=C["DGREY"], align="center", margin_in=0)
            l_start = box_y + 0.70
            l_space = (box_h - 0.75) / b["lines"]
            add_wlines(slide, b["x"] + 0.22, l_start + l_space * 0.9,
                       box_w - 0.44, b["lines"], l_space, C["PURPLE"])

    def build_slide2():
        slide = prs.slides.add_slide(BLANK)
        add_header(slide, "Word Maps",
                   "Read each word, break it into syllables, write the sound buttons, then write the word again.",
                   f"{CODE}.9")
        map_words = lesson["wordMaps"]["words"]
        hdr_row = [
            {"text": "Word",            "bold": True, "fill": C["LPURPLE"], "align": "center", "size": 11, "color": C["BLACK"]},
            {"text": "Syllable Breaks", "bold": True, "fill": C["LPURPLE"], "align": "center", "size": 11, "color": C["BLACK"]},
            {"text": "Sound Buttons",   "bold": True, "fill": C["LPURPLE"], "align": "center", "size": 11, "color": C["BLACK"]},
            {"text": "My Word",         "bold": True, "fill": C["LPURPLE"], "align": "center", "size": 11, "color": C["BLACK"]},
        ]
        data_rows = [
            [{"text": w, "bold": True, "fill": C["WHITE"], "align": "center", "size": 12, "color": C["BLACK"]},
             {"text": "", "fill": C["WHITE"]},
             {"text": "", "fill": C["WHITE"]},
             {"text": "", "fill": C["WHITE"]}]
            for w in map_words
        ]
        rh = [0.36] + [0.676] * len(data_rows)
        add_table(slide, [hdr_row] + data_rows,
                  MARGIN, CONTENT_Y + 0.04, CW, CONTENT_H - 0.08,
                  col_widths=[1.55, 2.10, 3.40, 2.25], row_heights=rh,
                  border_hex=C["BDR"])

    def build_slide3():
        slide = prs.slides.add_slide(BLANK)
        add_header(slide, "Definitions and In a Sentence",
                   "Read each word and its definition. Then write your own sentence using that word.",
                   f"{CODE}.11")
        hdr_row = [
            {"text": "Word",          "bold": True, "fill": C["LPURPLE"], "align": "center", "size": 10, "color": C["BLACK"]},
            {"text": "Definition",    "bold": True, "fill": C["LPURPLE"], "align": "left",   "size": 10, "color": C["BLACK"]},
            {"text": "In a Sentence", "bold": True, "fill": C["LPURPLE"], "align": "center", "size": 10, "color": C["BLACK"]},
        ]
        data_rows = [
            [{"text": w,       "bold": True,  "fill": C["WHITE"], "align": "center", "size": 11, "color": C["BLACK"]},
             {"text": DEFS[w], "bold": False, "fill": C["WHITE"], "align": "left",   "size": 8,  "color": C["DGREY"]},
             {"text": "",      "fill": C["WHITE"]}]
            for w in WORDS
        ]
        rh = [0.35] + [0.406] * len(data_rows)
        add_table(slide, [hdr_row] + data_rows,
                  MARGIN, CONTENT_Y + 0.04, CW, CONTENT_H - 0.08,
                  col_widths=[1.30, 3.45, 4.55], row_heights=rh,
                  border_hex=C["BDR"])

    def build_slide4():
        slide = prs.slides.add_slide(BLANK)
        add_header(slide, "Spell Check",
                   "Circle the correct spelling of each word.", f"{CODE}.13")
        rows = [
            [{"text": opt, "align": "center", "fill": C["WHITE"], "size": 14, "color": C["BLACK"]}
             for opt in row_data["opts"]]
            for row_data in SPELL_DATA
        ]
        add_table(slide, rows,
                  MARGIN, CONTENT_Y + 0.04, CW, CONTENT_H - 0.08,
                  col_widths=[3.1, 3.1, 3.1], row_heights=0.441,
                  border_hex=C["BDR"])

    def build_slide5():
        slide = prs.slides.add_slide(BLANK)
        add_header(slide, "Cloze",
                   "Use the words from the word bank to complete each sentence.", f"{CODE}.15")
        row_h = (CONTENT_H - 0.06) / len(WORDS)
        wb_w  = 1.25
        BLANK_STR = "____________________________"
        wb_rows = [[{"text": w, "align": "center", "bold": True,
                     "fill": C["WHITE"], "size": 10, "color": C["BLACK"], "valign": "middle"}]
                   for w in WORDS]
        add_table(slide, wb_rows,
                  MARGIN, CONTENT_Y + 0.03, wb_w, CONTENT_H - 0.06,
                  col_widths=[wb_w], row_heights=row_h, border_hex=C["BDR"])
        sent_x = MARGIN + wb_w + 0.18
        sent_w = CW - wb_w - 0.18
        for i, w in enumerate(CLOZE_ORDER):
            sentence = SENTENCES[w]
            idx = sentence.find(w)
            if idx == -1:
                before, after = sentence, ""
            else:
                before = sentence[:idx]
                after  = sentence[idx + len(w):]
            runs = [
                (before,    {"size": 11, "color": C["BLACK"]}),
                (BLANK_STR, {"size": 11, "color": C["BLACK"]}),
                (after,     {"size": 11, "color": C["BLACK"]}),
            ]
            add_rich_text(slide, runs,
                          sent_x, CONTENT_Y + 0.03 + i * row_h, sent_w, row_h,
                          align="left", valign="middle", margin_in=0.03)
            if i < len(WORDS) - 1:
                add_rect(slide, sent_x,
                         CONTENT_Y + 0.03 + (i + 1) * row_h - 0.012,
                         sent_w, 0.010, "E8E8E8")

    def build_slide6():
        slide = prs.slides.add_slide(BLANK)
        add_header(slide, "Morphology Matrix",
                   "How many new words can you make by adding suffix(es) or prefix(es) to the base word?",
                   f"{CODE}.19")
        mm = lesson["morphMatrix"]
        base_w, base_x = 2.90, MARGIN
        base_y, base_h = CONTENT_Y + 0.08, CONTENT_H - 0.12
        add_rect(slide, base_x, base_y, base_w, base_h, C["WHITE"],
                 line_hex=C["GREEN"], line_pt=2, radius=0.10)
        add_text(slide, "Base Word",
                 base_x, base_y + 0.08, base_w, 0.30,
                 size=13, bold=True, color=C["GREEN"], align="center", margin_in=0)
        add_text(slide, mm["baseWord"],
                 base_x, base_y + 0.44, base_w, 1.80,
                 size=58, color=C["BLACK"], align="center", valign="middle", margin_in=0)
        add_text(slide, mm["def"],
                 base_x + 0.15, base_y + base_h - 0.55, base_w - 0.30, 0.50,
                 size=9, italic=True, color=C["GREY"], align="center", valign="middle", margin_in=0)
        grid_x = MARGIN + base_w + 0.35
        grid_w = SW - grid_x - MARGIN
        gcw    = (grid_w - 0.15) / 2
        grh    = (base_h - 0.38 - 0.10) / 3
        add_text(slide, "Add a suffix or prefix:",
                 grid_x, base_y + 0.04, grid_w, 0.30,
                 size=13, bold=True, color=C["PURPLE"], align="center", margin_in=0)
        for i, sf in enumerate(mm["suffixes"]):
            col = i % 2
            row = i // 2
            cx = grid_x + col * (gcw + 0.15)
            cy = base_y + 0.38 + row * (grh + 0.10)
            add_rect(slide, cx, cy, gcw, grh, C["LGREY"],
                     line_hex=C["PURPLE"], line_pt=1.5, radius=0.08)
            add_text(slide, sf, cx, cy + 0.06, gcw, 0.26,
                     size=14, bold=True, color=C["DGREY"], align="center", margin_in=0)
            add_rect(slide, cx + 0.18, cy + grh - 0.34, gcw - 0.36, 0.013, C["PURPLE"])

    def build_slide7():
        slide = prs.slides.add_slide(BLANK)
        add_header(slide, "Scrambled Syllables",
                   "These words have been scrambled into their syllables. Can you put them back together?",
                   "")
        syl_breaks = lesson.get("syllableBreaks", {})
        map_syl    = lesson["wordMaps"].get("syllables", {})
        col_w2 = (CW - 0.25) / 2
        row_h2 = CONTENT_H / 5
        for i, w in enumerate(WORDS):
            raw   = syl_breaks.get(w) or map_syl.get(w) or w
            parts = list(reversed(raw.split(" | ")))
            col   = i % 2
            row   = i // 2
            cx    = MARGIN + col * (col_w2 + 0.25)
            cy    = CONTENT_Y + 0.02 + row * row_h2
            runs  = []
            for pi, part in enumerate(parts):
                runs.append((part, {"size": 16, "color": C["BLACK"]}))
                if pi < len(parts) - 1:
                    runs.append((" | ", {"size": 16, "color": C["RED"], "bold": True}))
            add_rich_text(slide, runs, cx, cy + 0.04, col_w2, 0.38,
                          align="left", valign="middle", margin_in=0)
            add_rect(slide, cx, cy + row_h2 - 0.18, col_w2 * 0.78, 0.013, C["PURPLE"])

    # Build all slides
    build_slide1()
    build_slide2()
    build_slide3()
    build_slide4()
    build_slide5()
    build_slide6()
    build_slide7()

    # Return as bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
