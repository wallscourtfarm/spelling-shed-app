"""
slides_builder.py
python-pptx port of slides-template.js — teaching deck.
Y4/Y5/Y6: 22 slides. Y3: 18 slides. Y2: 16 slides (+ optional morph matrix).
Returns raw PPTX bytes.
"""

import io
import os
import random
import zipfile
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


def build_slides(lesson: dict) -> bytes:

    WORDS       = lesson["words"]
    DEFS        = lesson["defs"]
    SENTENCES   = lesson["sentences"]
    CLOZE_ORDER = lesson["clozeOrder"]
    CODE        = lesson["code"]
    YEAR_GROUP  = lesson.get("yearGroup", "Y4")

    def rgb(h):
        return RGBColor.from_string(h)

    C = dict(
        BLUE    = "1E5FA3",
        PINK    = "E91E8C",
        PURPLE  = "7B1FA2",
        GREEN_S = "388E3C",
        CYAN    = "4DD0E1",
        GREEN_B = "57A657",
        BLACK   = "1A1A1A",
        GREY    = "757575",
        WHITE   = "FFFFFF",
        YELLOW  = "F4C430",
        RED     = "C62828",
    )

    SW, SH   = 10.0, 5.625
    HDR_H    = 1.12
    CYAN_H   = 0.06
    CONT_Y   = HDR_H + CYAN_H
    CONT_H   = 4.22
    BAR_Y    = CONT_Y + CONT_H
    FONT     = "Calibri"

    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    BLANK = prs.slide_layouts[6]

    slide_animations = []

    def register_clicks(slide, *click_groups):
        groups = []
        for cg in click_groups:
            if cg is None:
                continue
            if not isinstance(cg, (list, tuple)):
                cg = [cg]
            cg = [s for s in cg if s is not None]
            if cg:
                groups.append(cg)
        if groups:
            slide_animations.append((slide, groups))

    def inject_click_animations():
        P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        nsmap = {"p": P}

        def Pe(parent, tag, **attrs):
            el = etree.SubElement(parent, "{%s}%s" % (P, tag))
            for k, v in attrs.items():
                el.set(k, str(v))
            return el

        for slide, click_groups in slide_animations:
            sld = slide._element
            for tag in ("timing", "bldLst"):
                existing = sld.find("{%s}%s" % (P, tag))
                if existing is not None:
                    sld.remove(existing)

            timing = etree.SubElement(sld, "{%s}timing" % P, nsmap=nsmap)
            tnLst = Pe(timing, "tnLst")
            root_par = Pe(tnLst, "par")
            root_cTn = Pe(root_par, "cTn", id="1", dur="indefinite",
                          restart="whenNotActive", nodeType="tmRoot")
            root_children = Pe(root_cTn, "childTnLst")
            seq = Pe(root_children, "seq", concurrent="1", nextAc="seek")
            seq_cTn = Pe(seq, "cTn", id="2", dur="indefinite", nodeType="mainSeq")
            seq_children = Pe(seq_cTn, "childTnLst")

            cid = 3
            all_spids = []

            for group in click_groups:
                click_par = Pe(seq_children, "par")
                click_cTn = Pe(click_par, "cTn", id=str(cid), fill="hold")
                cid += 1
                stCondLst = Pe(click_cTn, "stCondLst")
                Pe(stCondLst, "cond", delay="indefinite")
                click_children = Pe(click_cTn, "childTnLst")

                inner_par = Pe(click_children, "par")
                inner_cTn = Pe(inner_par, "cTn", id=str(cid), fill="hold")
                cid += 1
                inner_stCond = Pe(inner_cTn, "stCondLst")
                Pe(inner_stCond, "cond", delay="0")
                inner_children = Pe(inner_cTn, "childTnLst")

                for shape_index, shape in enumerate(group):
                    spid = shape.shape_id
                    all_spids.append(spid)
                    node_type = "clickEffect" if shape_index == 0 else "withEffect"

                    effect_par = Pe(inner_children, "par")
                    effect_cTn = Pe(effect_par, "cTn", id=str(cid),
                                    presetID="1", presetClass="entr",
                                    presetSubtype="0", fill="hold",
                                    grpId="0", nodeType=node_type)
                    cid += 1
                    es_stCond = Pe(effect_cTn, "stCondLst")
                    Pe(es_stCond, "cond", delay="0")
                    e_children = Pe(effect_cTn, "childTnLst")

                    set_el = Pe(e_children, "set")
                    cBhvr = Pe(set_el, "cBhvr")
                    bhvr_cTn = Pe(cBhvr, "cTn", id=str(cid), dur="1", fill="hold")
                    cid += 1
                    bhvr_stCond = Pe(bhvr_cTn, "stCondLst")
                    Pe(bhvr_stCond, "cond", delay="0")
                    tgtEl = Pe(cBhvr, "tgtEl")
                    Pe(tgtEl, "spTgt", spid=str(spid))
                    attrNameLst = Pe(cBhvr, "attrNameLst")
                    attrName = etree.SubElement(attrNameLst, "{%s}attrName" % P)
                    attrName.text = "style.visibility"
                    to_el = Pe(set_el, "to")
                    Pe(to_el, "strVal", val="visible")

            bldLst = etree.SubElement(sld, "{%s}bldLst" % P, nsmap=nsmap)
            for spid in all_spids:
                Pe(bldLst, "bldP", spid=str(spid), grpId="0", animBg="1")

            prevCondLst = Pe(seq, "prevCondLst")
            prev_cond = Pe(prevCondLst, "cond", evt="onPrev", delay="0")
            prev_tgt = Pe(prev_cond, "tgtEl")
            Pe(prev_tgt, "sldTgt")
            nextCondLst = Pe(seq, "nextCondLst")
            next_cond = Pe(nextCondLst, "cond", evt="onNext", delay="0")
            next_tgt = Pe(next_cond, "tgtEl")
            Pe(next_tgt, "sldTgt")

    # ── Core drawing helpers ──────────────────────────────────────────────────

    def _bodyPr(txBody, valign="middle", margin_in=0.05):
        bp = txBody.find(qn('a:bodyPr'))
        if bp is None:
            bp = etree.SubElement(txBody, qn('a:bodyPr'))
        bp.set('anchor', {"top": "t", "middle": "ctr", "bottom": "b"}.get(valign, "ctr"))
        m = str(Inches(margin_in))
        for a in ('lIns', 'rIns', 'tIns', 'bIns'):
            bp.set(a, m)

    def rect(slide, x, y, w, h, fill, line=None, lpt=0, radius=None):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if line and lpt > 0:
            shape.line.color.rgb = rgb(line)
            shape.line.width = Pt(lpt)
        else:
            shape.line.fill.background()
        if radius is not None:
            sp = shape._element
            spPr = sp.find(qn('p:spPr'))
            pg = spPr.find(qn('a:prstGeom'))
            if pg is None:
                pg = etree.SubElement(spPr, qn('a:prstGeom'))
            pg.set('prst', 'roundRect')
            al = pg.find(qn('a:avLst'))
            if al is None:
                al = etree.SubElement(pg, qn('a:avLst'))
            for av in list(al):
                al.remove(av)
            adj = int(min(radius / min(w, h), 0.5) * 100000)
            gd = etree.SubElement(al, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {adj}')
        return shape

    def triangle(slide, x, y, w, h, fill, line=None, lpt=1):
        shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if line and lpt > 0:
            shape.line.color.rgb = rgb(line)
            shape.line.width = Pt(lpt)
        else:
            shape.line.fill.background()
        return shape

    def txt(slide, text, x, y, w, h, *,
            font=FONT, size=14, bold=False, italic=False,
            color="1A1A1A", align="left", valign="middle",
            wrap=True, margin=0.04, shadow=False):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        _bodyPr(tf._txBody, valign=valign, margin_in=margin)
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = text
        rf = run.font
        rf.name = font
        rf.size = Pt(size)
        rf.bold = bold
        rf.italic = italic
        rf.color.rgb = rgb(color)
        return tb

    def rich_txt(slide, runs, x, y, w, h, *,
                 font=FONT, align="left", valign="middle", margin=0.04):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        _bodyPr(tf._txBody, valign=valign, margin_in=margin)
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER}.get(align, PP_ALIGN.LEFT)
        for text, opts in runs:
            run = p.add_run()
            run.text = text
            rf = run.font
            rf.name = font
            rf.size = Pt(opts.get("size", 14))
            rf.bold = opts.get("bold", False)
            rf.italic = opts.get("italic", False)
            rf.color.rgb = rgb(opts.get("color", "1A1A1A"))
        return tb

    def table(slide, rows, x, y, w, h, col_widths, row_heights,
              border_hex="444444", border_pt=1.5, font=FONT):
        nr, nc = len(rows), len(rows[0])
        ts = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
        tbl = ts.table
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)
        rhl = row_heights if isinstance(row_heights, list) else [row_heights] * nr
        for ri, rh in enumerate(rhl):
            tbl.rows[ri].height = Inches(rh)
        for ri, row in enumerate(rows):
            for ci, cd in enumerate(row):
                cell = tbl.cell(ri, ci)
                text   = cd.get("text", "")
                bold   = cd.get("bold", False)
                italic = cd.get("italic", False)
                color  = cd.get("color", "1A1A1A")
                fill   = cd.get("fill", "FFFFFF")
                align  = cd.get("align", "left")
                size   = cd.get("size", 14)
                valign = cd.get("valign", "middle")
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(fill)
                tf = cell.text_frame
                tf.word_wrap = True
                _bodyPr(tf._txBody, valign=valign, margin_in=0.04)
                p = tf.paragraphs[0]
                p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                               "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
                if text:
                    run = p.add_run()
                    run.text = text
                    rf = run.font
                    rf.name = font
                    rf.size = Pt(size)
                    rf.bold = bold
                    rf.italic = italic
                    rf.color.rgb = rgb(color)
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
                    ln = tcPr.find(qn(f'a:{edge}'))
                    if ln is None:
                        ln = etree.SubElement(tcPr, qn(f'a:{edge}'))
                    ln.set('w', str(int(Pt(border_pt))))
                    sf = ln.find(qn('a:solidFill'))
                    if sf is None:
                        sf = etree.SubElement(ln, qn('a:solidFill'))
                    sc = sf.find(qn('a:srgbClr'))
                    if sc is None:
                        sc = etree.SubElement(sf, qn('a:srgbClr'))
                    sc.set('val', border_hex)
        return ts

    # ── Frame ─────────────────────────────────────────────────────────────────

    def add_frame(slide, activity_type, activity_label, question, slide_num):
        ind_col = "F4C430" if activity_type == "Independent" else "66BB6A"
        rect(slide, 0.1, 0.1, 1.35, 0.92, ind_col)
        txt(slide, activity_type, 0.1, 0.1, 1.35, 0.92,
            size=12, bold=True, color=C["WHITE"], align="center", margin=0)
        txt(slide, activity_label, 1.55, 0.06, 7.5, 0.44,
            size=19, color=C["BLACK"], align="center")
        txt(slide, question, 1.55, 0.5, 7.5, 0.6,
            size=22, bold=True, color=C["BLUE"], align="center")
        txt(slide, slide_num, 8.85, 0.78, 1.05, 0.28,
            size=14, color=C["GREY"], align="right")
        rect(slide, 0, HDR_H, SW, CYAN_H, C["CYAN"])
        rect(slide, 0, BAR_Y, SW, SH - BAR_Y, C["GREEN_B"])

    def fit_font(text, box_w, max_pt=36, min_pt=14):
        n = len(text.replace(" ", "")) or len(text)
        req = (box_w * 72) / (n * 0.54)
        return round(min(max_pt, max(min_pt, req)))

    # ── Sound button drawing ──────────────────────────────────────────────────

    def _draw_bezier_arc(slide, x1, y1, x2, y2, peak_y, line_w):
        EMU = 914400
        cx = (x1 + x2) / 2
        mid_y = (y1 + y2) / 2
        cy = 2 * peak_y - mid_y

        bx_min = min(x1, x2, cx)
        bx_max = max(x1, x2, cx)
        by_min = min(y1, y2, peak_y, cy)
        by_max = max(y1, y2, peak_y, cy)
        box_x = bx_min
        box_y = by_min
        box_w = max(bx_max - bx_min, 0.001)
        box_h = max(by_max - by_min, 0.001)

        path_max = 21600
        def to_path_x(x):
            return int((x - box_x) / box_w * path_max) if box_w > 0 else 0
        def to_path_y(y):
            return int((y - box_y) / box_h * path_max) if box_h > 0 else 0

        sp = etree.SubElement(slide.shapes._spTree, qn('p:sp'))
        nvSpPr = etree.SubElement(sp, qn('p:nvSpPr'))
        cNvPr = etree.SubElement(nvSpPr, qn('p:cNvPr'))
        cNvPr.set('id', str(slide.shapes._next_shape_id))
        cNvPr.set('name', 'BezierArc')
        etree.SubElement(nvSpPr, qn('p:cNvSpPr'))
        etree.SubElement(nvSpPr, qn('p:nvPr'))
        spPr = etree.SubElement(sp, qn('p:spPr'))

        xfrm = etree.SubElement(spPr, qn('a:xfrm'))
        off = etree.SubElement(xfrm, qn('a:off'))
        off.set('x', str(int(box_x * EMU)))
        off.set('y', str(int(box_y * EMU)))
        ext = etree.SubElement(xfrm, qn('a:ext'))
        ext.set('cx', str(int(box_w * EMU)))
        ext.set('cy', str(int(box_h * EMU)))

        custGeom = etree.SubElement(spPr, qn('a:custGeom'))
        etree.SubElement(custGeom, qn('a:avLst'))
        etree.SubElement(custGeom, qn('a:gdLst'))
        etree.SubElement(custGeom, qn('a:ahLst'))
        etree.SubElement(custGeom, qn('a:cxnLst'))
        etree.SubElement(custGeom, qn('a:rect')).attrib.update(
            {'l': '0', 't': '0', 'r': str(path_max), 'b': str(path_max)})

        pathLst = etree.SubElement(custGeom, qn('a:pathLst'))
        path = etree.SubElement(pathLst, qn('a:path'))
        path.set('w', str(path_max))
        path.set('h', str(path_max))

        moveTo = etree.SubElement(path, qn('a:moveTo'))
        pt = etree.SubElement(moveTo, qn('a:pt'))
        pt.set('x', str(to_path_x(x1)))
        pt.set('y', str(to_path_y(y1)))

        quadBezTo = etree.SubElement(path, qn('a:quadBezTo'))
        ctrl = etree.SubElement(quadBezTo, qn('a:pt'))
        ctrl.set('x', str(to_path_x(cx)))
        ctrl.set('y', str(to_path_y(cy)))
        end = etree.SubElement(quadBezTo, qn('a:pt'))
        end.set('x', str(to_path_x(x2)))
        end.set('y', str(to_path_y(y2)))

        noFill = etree.SubElement(spPr, qn('a:noFill'))
        ln = etree.SubElement(spPr, qn('a:ln'))
        ln.set('w', str(int(line_w * 12700)))
        ln.set('cap', 'rnd')
        solidFill = etree.SubElement(ln, qn('a:solidFill'))
        clr = etree.SubElement(solidFill, qn('a:srgbClr'))
        clr.set('val', C["BLACK"])

    def draw_sound_buttons(slide, word, cx, wy, font_size):
        phonemes = lesson["phonemes"].get(word)
        if not phonemes:
            return

        CELL_W = max(font_size / 28 * 0.40, font_size * 0.018)
        ROW_H  = font_size * 1.4 / 72
        DOT    = max(0.040, font_size * 0.084 / 28)
        LINE_H = max(0.018, font_size * 0.046 / 28)
        PAD    = CELL_W * 0.07
        GAP    = 0.04

        n_cells = sum(len(g["l"]) for g in phonemes)
        total_w = n_cells * CELL_W
        table_x = cx - total_w / 2
        sym_y   = wy + ROW_H + GAP
        line_y  = sym_y + (DOT - LINE_H) / 2

        ts = slide.shapes.add_table(
            1, n_cells,
            Inches(table_x), Inches(wy),
            Inches(total_w), Inches(ROW_H)
        )
        tbl = ts.table
        for ci in range(n_cells):
            tbl.columns[ci].width = Inches(CELL_W)
        tbl.rows[0].height = Inches(ROW_H)

        ci = 0
        for g in phonemes:
            for ch in g["l"]:
                cell = tbl.cell(0, ci)
                cell.fill.background()
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                tcPr.set('marL', '0')
                tcPr.set('marR', '0')
                tcPr.set('marT', '0')
                tcPr.set('marB', '0')
                tf = cell.text_frame
                _bodyPr(tf._txBody, valign="middle", margin_in=0)
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = ch
                run.font.name = FONT
                run.font.size = Pt(int(font_size))
                run.font.bold = False
                run.font.color.rgb = rgb(C["BLACK"])
                for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
                    ln = tcPr.find(qn(f'a:{edge}'))
                    if ln is None:
                        ln = etree.SubElement(tcPr, qn(f'a:{edge}'))
                    for sf in ln.findall(qn('a:solidFill')):
                        ln.remove(sf)
                    if ln.find(qn('a:noFill')) is None:
                        etree.SubElement(ln, qn('a:noFill'))
                ci += 1

        col = 0
        gd = []
        for g in phonemes:
            n  = len(g["l"])
            gx = table_x + col * CELL_W
            gw = n * CELL_W
            gd.append({**g, "gx": gx, "gw": gw, "gc": gx + gw / 2})
            col += n

        for g in gd:
            if g["t"] == "line":
                rect(slide, g["gx"] + PAD, line_y,
                     g["gw"] - 2 * PAD, LINE_H, C["BLACK"])

        for g in gd:
            if g["t"] == "dot" and "sid" not in g:
                shape = slide.shapes.add_shape(
                    9,
                    Inches(g["gc"] - DOT / 2), Inches(sym_y),
                    Inches(DOT), Inches(DOT)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb(C["BLACK"])
                shape.line.fill.background()

        splits = {}
        for g in gd:
            if "sid" in g:
                splits.setdefault(g["sid"], []).append(g)

        for sid, pair in splits.items():
            if len(pair) < 2:
                continue
            pair.sort(key=lambda g: g["gx"])
            x1, x2 = pair[0]["gc"], pair[-1]["gc"]
            arc_h = min(0.14, max(0.08, font_size * 0.005))
            base_y = sym_y + DOT * 0.5
            peak_y = base_y + arc_h
            line_w = max(1.0, font_size * 0.05)
            _draw_bezier_arc(slide, x1, base_y, x2, base_y, peak_y, line_w)

    def draw_phoneme_symbols_only(slide, word, cx, wy, scale=1.0):
        """Draw phoneme dots/lines/arcs without the letter table. Returns total width drawn."""
        phonemes = lesson["phonemes"].get(word)
        if not phonemes:
            return 0.0
        UNIT   = max(0.12, 0.15 * scale)
        DOT    = max(0.055, 0.065 * scale)
        LINE_H = max(0.015, 0.018 * scale)
        PAD    = max(0.008, 0.01 * scale)
        n_cells = sum(len(g["l"]) for g in phonemes)
        total_w = n_cells * UNIT
        sx = cx - total_w / 2
        col = 0
        gd = []
        for g in phonemes:
            n = len(g["l"])
            gx = sx + col * UNIT
            gw = n * UNIT
            gd.append({**g, "gx": gx, "gw": gw, "gc": gx + gw / 2})
            col += n
        line_y = wy + (DOT - LINE_H) / 2
        for g in gd:
            if g["t"] == "line":
                rect(slide, g["gx"] + PAD, line_y,
                     max(g["gw"] - 2 * PAD, 0.001), LINE_H, C["BLACK"])
        for g in gd:
            if g["t"] == "dot" and "sid" not in g:
                sh = slide.shapes.add_shape(
                    9, Inches(g["gc"] - DOT / 2), Inches(wy), Inches(DOT), Inches(DOT))
                sh.fill.solid()
                sh.fill.fore_color.rgb = rgb(C["BLACK"])
                sh.line.fill.background()
        splits = {}
        for g in gd:
            if "sid" in g:
                splits.setdefault(g["sid"], []).append(g)
        for sid, pair in splits.items():
            if len(pair) < 2:
                continue
            pair.sort(key=lambda g: g["gx"])
            x1, x2 = pair[0]["gc"], pair[-1]["gc"]
            arc_h = max(0.04, 0.06 * scale)
            base_y = wy + DOT * 0.5
            _draw_bezier_arc(slide, x1, base_y, x2, base_y, base_y + arc_h, max(0.8, 1.5 * scale))
        return total_w

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title (shared all year groups)
    # ════════════════════════════════════════════════════════════════════════

    def slide1():
        s = prs.slides.add_slide(BLANK)
        rect(s, 0, 0, SW, 3.6, "87CEEB")
        rect(s, 0, 3.3, SW, 2.325, C["WHITE"])
        rect(s, 0, BAR_Y, SW, SH - BAR_Y, C["GREEN_B"])
        txt(s, "Spelling Shed", 1.5, 0.15, 7.0, 1.4,
            size=64, bold=True, color=C["YELLOW"], align="center")

        def badge(slide, text, x, y, w, h):
            rect(slide, x, y, w, h, C["WHITE"], "BBBBBB", 1.5, radius=0.15)
            txt(slide, text, x, y, w, h, size=20, bold=True,
                color=C["BLACK"], align="center", margin=0)

        stage_val = str(lesson["stage"]).strip()
        if stage_val.lower().startswith("stage"):
            stage_val = stage_val[5:].strip().lstrip(":").strip()
        badge(s, "Stage: " + stage_val, 3.2, 1.75, 1.6, 0.75)
        badge(s, "Lesson: " + CODE,     5.2, 1.75, 1.6, 0.75)

        obj = f"To spell words: {lesson['rule']}"
        txt(s, obj, 0.5, 3.42, 9.0, 0.42,
            size=18, bold=True, color=C["BLACK"], align="center")
        txt(s, "This week's words:  " + ", ".join(WORDS),
            0.5, 4.35, 9.0, 0.55, size=17, color=C["BLACK"], align="center")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE — Key Spelling Practice (shared all year groups)
    # ════════════════════════════════════════════════════════════════════════

    key_spelling_slide_index = [None]
    key_spelling_pic_id = [None]

    def slide_key_spelling():
        if not lesson.get("keySpellingWord"):
            return
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Key Spelling Practice",
                  "Quick Write word", "KSp")

        rect(s, 0.58, 1.42, 8.23, 0.45, "FFF9C4", C["YELLOW"], 1.5)
        txt(s, "Today's word is…", 0.58, 1.42, 8.23, 0.45,
            size=16, bold=True, color=C["BLACK"], align="center", margin=0)

        word = str(lesson["keySpellingWord"])
        word_fs = 66
        if len(word) > 9:
            word_fs = max(40, int(66 * 9 / len(word)))
        txt(s, word, 1.0, 2.34, 8.0, 1.21,
            size=word_fs, color=C["BLACK"], align="center", valign="middle", margin=0)

        txt(s, "How many times can you write the word in 1 minute?",
            2.0, 4.27, 6.0, 0.40,
            size=14, color=C["BLACK"], align="center", margin=0)

        timer_img = os.path.join(os.path.dirname(__file__), "assets", "timer_1min_thumb.png")
        if os.path.exists(timer_img):
            pic = s.shapes.add_picture(
                timer_img,
                Inches(4.38), Inches(4.83),
                Inches(1.40), Inches(0.42)
            )
            key_spelling_slide_index[0] = len(prs.slides) - 1
            key_spelling_pic_id[0] = pic.shape_id

    # ════════════════════════════════════════════════════════════════════════
    # Y4/Y5/Y6 SLIDES
    # ════════════════════════════════════════════════════════════════════════

    def slide2():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Starter",
                  lesson["starter"]["question"], f"{CODE}.1")
        txt(s, "The root or base word may need to change.",
            1.0, CONT_Y + 0.15, 8.0, 0.4,
            size=18, color=C["BLACK"], align="center")

        starters = lesson["starter"]["words"]
        cell_w, cell_h = 9.0 / 3, 1.5
        sx, sy = 0.5, CONT_Y + 0.65
        label = lesson["starter"]["answerLabel"]

        for i, word in enumerate(starters):
            col, row = i % 3, i // 3
            x = sx + col * cell_w
            y = sy + row * cell_h
            txt(s, word, x, y, cell_w, 0.8,
                size=44, color=C["BLACK"], align="center")
            txt(s, f"{label} ______________________",
                x + 0.15, y + 1.1, cell_w - 0.3, 0.3,
                size=14, color=C["BLACK"], align="left")

    def slide3():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Starter",
                  "How did the base word change?", f"{CODE}.2")

        starters = lesson["starter"]["words"]
        answers  = lesson["starter"]["answers"]
        cell_w   = 9.0 / 3
        sx, sy   = 0.5, CONT_Y + 0.15
        BASE_H, ANS_H, RULE_H = 0.52, 0.50, 0.28
        PAIR_H   = BASE_H + ANS_H + RULE_H
        row_gap  = 0.20

        fs = max(14, min(28, int((cell_w * 72) / (max(len(w) for w in answers) * 0.54))))

        click_groups = []
        for i, (word, ans) in enumerate(zip(starters, answers)):
            col, row = i % 3, i // 3
            x = sx + col * cell_w
            y = sy + row * (PAIR_H + row_gap)
            txt(s, word + "  →", x, y, cell_w, BASE_H,
                size=fs, color=C["BLACK"], align="center")
            ans_shape = txt(s, ans, x, y + BASE_H, cell_w, ANS_H,
                size=fs, bold=True, color=C["PINK"], align="center")
            note = lesson["starter"].get("perPairNote", "")
            note_shape = txt(s, note, x, y + BASE_H + ANS_H, cell_w, RULE_H,
                size=11, italic=True, color=C["GREY"], align="center")
            click_groups.append([ans_shape, note_shape])

        rule_y = sy + 2 * (PAIR_H + row_gap) + 0.14
        rule_box_shape = rect(s, 0.5, rule_y, 9.0, 0.9, "FFF9C4", C["YELLOW"], 1.5)
        rule_title_shape = txt(s, lesson["starter"]["ruleBox"], 0.5, rule_y + 0.04, 9.0, 0.32,
            size=15, bold=True, color=C["BLACK"], align="center", margin=0)
        rule_text = lesson["starter"].get("ruleText", "")
        rule_text_shape = None
        if rule_text:
            rule_text_shape = txt(s, rule_text, 0.7, rule_y + 0.36, 8.6, 0.5,
                size=12, color=C["BLACK"], align="center", margin=0)
        click_groups.append([rule_box_shape, rule_title_shape, rule_text_shape])

        register_clicks(s, *click_groups)

    def slide4():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "This Week's Words",
                  lesson["thisWeeksWordsQ"], f"{CODE}.3")

        cell_w = 9.0 / 5
        fs = min(fit_font(w, cell_w, 32, 18) for w in WORDS)

        for ri, row_words in enumerate([WORDS[:5], WORDS[5:]]):
            for ci, w in enumerate(row_words):
                x = 0.5 + ci * cell_w
                y = CONT_Y + 0.3 + ri * 1.30
                txt(s, w, x, y, cell_w, 0.85,
                    size=fs, color=C["BLACK"], align="center", valign="bottom")
                bar_w = cell_w - 0.6
                rect(s, x + 0.3, y + 0.92, bar_w, 0.06, C["PINK"])

        prompt_shape = txt(s, lesson["thisWeeksWordsPrompt"],
            0.5, CONT_Y + 3.05, 4.3, 0.4,
            size=15, bold=True, color=C["BLUE"], align="left")
        explanation_shape = txt(s, lesson["thisWeeksWordsExplanation"],
            4.8, CONT_Y + 3.05, 4.8, 0.55,
            size=14, color=C["BLACK"], align="left")
        register_clicks(s, prompt_shape, explanation_shape)

    def slide5():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Etymology",
                  "Which of this week's words is this?", f"{CODE}.4")

        etym = lesson["etymology"]
        word_shape = txt(s, etym["word"], 1.0, CONT_Y + 0.1, 8.0, 0.85,
            size=60, color=C["PINK"], align="center")

        card_groups = []
        bw = 2.85
        for i, click in enumerate(etym["clicks"][:3]):
            x = 0.5 + i * (bw + 0.15)
            box_shape = rect(s, x, CONT_Y + 1.1, bw, 1.35, C["WHITE"], C["PINK"], 2)
            text_shape = txt(s, click["label"] + " — " + click["body"],
                x + 0.12, CONT_Y + 1.15, bw - 0.24, 1.25,
                size=13, color=C["BLACK"], align="center")
            card_groups.append([box_shape, text_shape])

        base_word_shape = txt(s, etym["baseForm"], 1.25, CONT_Y + 2.55, 2.5, 0.40,
            size=26, color=C["BLUE"], align="center")
        arrows_shape = txt(s, "↙                          ↘",
            1.25, CONT_Y + 2.95, 2.5, 0.22,
            size=14, color=C["BLACK"], align="center")

        sub3_shape = sub4_shape = None
        if len(etym["clicks"]) > 3:
            sub3_shape = txt(s, etym["clicks"][3]["label"] + ":\n" + etym["clicks"][3]["body"],
                0.3, CONT_Y + 3.19, 2.1, 0.84,
                size=10, color=C["GREY"], align="center")
        if len(etym["clicks"]) > 4:
            sub4_shape = txt(s, etym["clicks"][4]["label"] + ":\n" + etym["clicks"][4]["body"],
                2.6, CONT_Y + 3.19, 2.4, 0.84,
                size=10, color=C["GREY"], align="center")

        right_label_shape = right_body_shape = None
        if len(etym["clicks"]) > 5:
            right_label_shape = txt(s, etym["clicks"][5]["label"] + ":",
                5.2, CONT_Y + 2.55, 4.3, 0.38,
                size=14, color=C["BLACK"], align="left")
            right_body_shape = txt(s, etym["clicks"][5]["body"],
                5.2, CONT_Y + 2.95, 4.3, 1.0,
                size=11, color=C["BLACK"], align="left", valign="top")

        register_clicks(s,
            *card_groups,
            word_shape,
            [base_word_shape, arrows_shape, sub3_shape, sub4_shape],
            [right_label_shape, right_body_shape],
        )

    def slide6():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Syllable Count",
                  "How many syllables are in this week's words?", f"{CODE}.5")

        cell_w, cell_h = 9.0 / 5, 1.6
        sx, sy = 0.5, CONT_Y + 0.2
        sc = lesson["syllableCounts"]
        fs = min(fit_font(w, cell_w, 28, 14) for w in WORDS)

        click_shapes = []
        for i, w in enumerate(WORDS):
            col, row = i % 5, i // 5
            x = sx + col * cell_w
            y = sy + row * cell_h
            n = sc.get(w, 1)
            txt(s, w, x, y, cell_w, 0.8,
                size=fs, color=C["BLACK"], align="center")
            num_col = C["RED"] if n == 1 else "388E3C"
            num_shape = txt(s, str(n), x, y + 0.82, cell_w, 0.42,
                size=26, bold=True, color=num_col, align="center")
            label_shape = txt(s, "syllable" if n == 1 else "syllables",
                x, y + 1.24, cell_w, 0.26,
                size=11, color=C["GREY"], align="center")
            click_shapes.extend([num_shape, label_shape])

        txt(s, "Does the length of the word affect the number of syllables it has?",
            0.5, BAR_Y - 0.48, 9.0, 0.35,
            size=13, italic=True, color=C["GREY"], align="center")

        register_clicks(s, click_shapes)

    def slide7():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Word Sort",
                  lesson["wordSortQ"], f"{CODE}.6")

        for ri, row_words in enumerate([WORDS[:5], WORDS[5:]]):
            txt(s, "          ".join(row_words),
                0.5, CONT_Y + 0.15 + ri * 0.5, 9.0, 0.48,
                size=22, color=C["BLACK"], align="center")

        BOX_Y, BOX_H = CONT_Y + 1.3, 2.05
        ws = lesson["wordSort"]

        rect(s, 0.35, BOX_Y, 4.45, BOX_H, C["WHITE"], C["PURPLE"], 3, radius=0.15)
        txt(s, ws["box1label"] + "\n" + ws["box1sub"],
            0.45, BOX_Y + 0.08, 4.25, 0.75,
            size=16, bold=True, color=C["PURPLE"], align="center")

        rect(s, 5.2, BOX_Y, 4.45, BOX_H, C["WHITE"], C["GREEN_S"], 3, radius=0.15)
        txt(s, ws["box2label"] + "\n" + ws["box2sub"],
            5.3, BOX_Y + 0.08, 4.25, 0.75,
            size=16, bold=True, color=C["GREEN_S"], align="center")

        txt(s, ws["hint"], 0.5, BOX_Y + BOX_H + 0.12, 9.0, 0.32,
            size=12, italic=True, color=C["GREY"], align="center")

    def slide8():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Word Sort", "Answers", f"{CODE}.7")

        ws = lesson["wordSort"]
        txt(s, ws["answerNote"], 0.5, CONT_Y + 0.08, 9.0, 0.38,
            size=14, color=C["BLACK"], align="center")

        BOX_Y, BOX_H = CONT_Y + 0.55, 3.05
        verb_only = ws["verbOnly"]
        verb_noun = ws["verbNoun"]

        rect(s, 0.35, BOX_Y, 4.45, BOX_H, C["WHITE"], C["PURPLE"], 3, radius=0.15)
        txt(s, ws["box1label"] + "  " + ws["box1sub"],
            0.45, BOX_Y + 0.08, 4.25, 0.35,
            size=14, bold=True, color=C["PURPLE"], align="center")

        if verb_only:
            available_h = BOX_H - 0.50 - 0.10
            row_h  = min(0.52, available_h / len(verb_only))
            font_s = max(9, min(13, int(row_h * 22)))
            for i, w in enumerate(verb_only):
                txt(s, w, 0.55, BOX_Y + 0.50 + i * row_h, 4.15, row_h,
                    size=font_s, bold=True, color=C["PINK"], align="center")

        rect(s, 5.2, BOX_Y, 4.45, BOX_H, C["WHITE"], C["GREEN_S"], 3, radius=0.15)
        txt(s, ws["box2label"], 5.3, BOX_Y + 0.08, 4.25, 0.35,
            size=14, bold=True, color=C["GREEN_S"], align="center")

        if verb_noun:
            available_h = BOX_H - 0.47 - 0.10
            row_h  = min(0.362, available_h / len(verb_noun))
            font_w = max(9, min(12, int(row_h * 28)))
            font_eg = max(8, font_w - 2)
            for i, item in enumerate(verb_noun):
                rich_txt(s, [
                    (item["word"], {"size": font_w, "bold": True,  "color": C["PINK"]}),
                    ("  " + item["eg"], {"size": font_eg, "italic": True, "color": C["GREY"]}),
                ], 5.4, BOX_Y + 0.47 + i * row_h, 4.15, row_h, valign="middle")

        txt(s, ws["exampleLine"], 0.5, BOX_Y + BOX_H + 0.12, 9.0, 0.30,
            size=12, color=C["BLACK"], align="center")

    def slide9():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Syllable and Phoneme Map",
                  "How to mark the syllable breaks and sound buttons", f"{CODE}.8")

        map_words = lesson["wordMaps"]["words"][:3]
        syllables = lesson["wordMaps"]["syllables"]
        sc        = lesson["syllableCounts"]

        cw, sx, sy = 2.9, 0.5, CONT_Y + 0.25
        fs = min(fit_font(w, cw, 36, 18) for w in map_words)

        for i, w in enumerate(map_words):
            x  = sx + i * (cw + 0.25)
            cx = x + cw / 2
            brk = syllables.get(w, w)
            n   = sc.get(w, 1)

            txt(s, w, x, sy, cw, 0.7,
                size=fs, color=C["PINK"], align="center")
            txt(s, "Syllable Breaks", x, sy + 0.75, cw, 0.3,
                size=12, bold=True, color=C["BLUE"], align="center")
            rect(s, x, sy + 1.08, cw, 0.5, C["WHITE"], C["BLUE"], 1.5)
            txt(s, brk, x, sy + 1.08, cw, 0.5,
                size=20, color=C["BLACK"], align="center", margin=0)
            txt(s, f"{n} syllable{'s' if n > 1 else ''}",
                x, sy + 1.62, cw, 0.28,
                size=12, italic=True, color=C["GREY"], align="center")
            txt(s, "Sound Buttons", x, sy + 1.95, cw, 0.3,
                size=12, bold=True, color=C["PINK"], align="center")
            draw_sound_buttons(s, w, cx, sy + 2.32, 22)
            txt(s, "● = single sound   — = letters making one sound",
                x, sy + 2.88, cw, 0.22,
                size=9, italic=True, color=C["GREY"], align="center")

    def slide10():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Word Maps",
                  "Read the word. Write syllable breaks and sound buttons. Rewrite the word.",
                  f"{CODE}.9")

        map_words = lesson["wordMaps"]["words"]
        hdr = [
            {"text": "Words",           "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "Syllable Breaks", "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "Sound Buttons",   "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "My Word",         "bold": True, "fill": C["WHITE"], "align": "center"},
        ]
        rows = [
            [{"text": w, "fill": C["WHITE"], "align": "center", "size": 18},
             {"text": "", "fill": C["WHITE"]},
             {"text": "", "fill": C["WHITE"]},
             {"text": "", "fill": C["WHITE"]}]
            for w in map_words
        ]
        rh = [0.5] + [0.57] * len(rows)
        table(s, [hdr] + rows,
              0.35, CONT_Y + 0.1, 9.3, CONT_H - 0.15,
              [2.0, 2.5, 2.9, 1.9], rh)

    def slide11():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Word Maps", "Answers", f"{CODE}.10")

        map_words = lesson["wordMaps"]["words"]
        syllables = lesson["wordMaps"]["syllables"]
        hdr = [
            {"text": "Words",           "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "Syllable Breaks", "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "Sound Buttons",   "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "My Word",         "bold": True, "fill": C["WHITE"], "align": "center"},
        ]
        rows = [
            [{"text": w,                   "fill": C["WHITE"], "align": "center", "size": 18},
             {"text": syllables.get(w, w), "fill": C["WHITE"], "align": "center", "size": 16},
             {"text": "",                  "fill": C["WHITE"]},
             {"text": w,                   "fill": C["WHITE"], "align": "center", "size": 18, "color": C["PINK"]}]
            for w in map_words
        ]
        rh = [0.5] + [0.57] * len(rows)
        ts = table(s, [hdr] + rows,
                   0.35, CONT_Y + 0.1, 9.3, CONT_H - 0.15,
                   [2.0, 2.5, 2.9, 1.9], rh)

        TABLE_X   = 0.35
        SB_COL_X  = TABLE_X + 2.0 + 2.5
        SB_COL_W  = 2.9
        HDR_H_IN  = 0.5
        ROW_H_IN  = 0.57
        SB_FS     = 13

        for ri, w in enumerate(map_words):
            cell_top_y = CONT_Y + 0.1 + HDR_H_IN + ri * ROW_H_IN
            cell_cx    = SB_COL_X + SB_COL_W / 2
            row_h_sb   = SB_FS * 1.4 / 72
            dot_h      = max(0.040, SB_FS * 0.084 / 28)
            total_sb_h = row_h_sb + 0.04 + dot_h
            word_top_y = cell_top_y + (ROW_H_IN - total_sb_h) / 2
            draw_sound_buttons(s, w, cell_cx, word_top_y, SB_FS)

    def slide12():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Definitions and In a Sentence",
                  "Read the word and its definition, then write your own sentence for each word.",
                  f"{CODE}.11")

        hdr = [
            {"text": "Word",          "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "Definition",    "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "In a Sentence", "bold": True, "fill": C["WHITE"], "align": "center"},
        ]
        rows = [
            [{"text": w,       "fill": C["WHITE"], "align": "center", "size": 13},
             {"text": DEFS[w], "fill": C["WHITE"], "align": "left",   "size": 9},
             {"text": "",      "fill": C["WHITE"]}]
            for w in WORDS
        ]
        rh = [0.36] + [0.34] * len(rows)
        table(s, [hdr] + rows,
              0.35, CONT_Y + 0.05, 9.3, CONT_H - 0.1,
              [1.4, 3.7, 4.2], rh)

    def slide13():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Definitions and In a Sentence",
                  "Answers", f"{CODE}.12")

        hdr = [
            {"text": "Word",          "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "Definition",    "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "In a Sentence", "bold": True, "fill": C["WHITE"], "align": "center"},
        ]
        rows = [
            [{"text": w,            "fill": C["WHITE"], "align": "center", "size": 13},
             {"text": DEFS[w],      "fill": C["WHITE"], "align": "left",   "size": 9},
             {"text": SENTENCES[w], "fill": C["WHITE"], "align": "left",   "size": 9,
              "color": C["PINK"]}]
            for w in WORDS
        ]
        rh = [0.36] + [0.34] * len(rows)
        table(s, [hdr] + rows,
              0.35, CONT_Y + 0.05, 9.3, CONT_H - 0.1,
              [1.4, 3.55, 4.35], rh)

    def slide14():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Spell Check",
                  "Circle the correct spelling of each word.", f"{CODE}.13")

        rows = [
            [{"text": opt, "fill": C["WHITE"], "align": "center", "size": 16}
             for opt in row["opts"]]
            for row in lesson["spellData"]
        ]
        table(s, rows,
              0.35, CONT_Y + 0.05, 9.3, CONT_H - 0.05,
              [3.1, 3.1, 3.1], 0.4)

    def slide15():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Spell Check", "Answers", f"{CODE}.14")

        rows = [
            [{"text": opt,
              "fill": C["WHITE"],
              "align": "center",
              "size": 16,
              "bold": j == row["correct"],
              "color": C["PINK"] if j == row["correct"] else C["BLACK"]}
             for j, opt in enumerate(row["opts"])]
            for row in lesson["spellData"]
        ]
        table(s, rows,
              0.35, CONT_Y + 0.05, 9.3, CONT_H - 0.05,
              [3.1, 3.1, 3.1], 0.4)

    def slide16():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Cloze",
                  "Use the words from the word bank to complete each sentence.",
                  f"{CODE}.15")

        row_h = (CONT_H - 0.10) / len(WORDS)
        wb_rows = [[{"text": w, "fill": C["WHITE"], "align": "center",
                     "size": 12, "valign": "middle"}]
                   for w in WORDS]
        table(s, wb_rows,
              0.30, CONT_Y + 0.05, 1.50, CONT_H - 0.10,
              [1.50], row_h)

        BLANK_STR = "_______________"
        for i, w in enumerate(CLOZE_ORDER):
            sentence = SENTENCES[w]
            idx = sentence.find(w)
            before = sentence[:idx] if idx != -1 else sentence
            after  = sentence[idx + len(w):] if idx != -1 else ""
            rich_txt(s, [
                (before,     {"size": 12, "color": C["BLACK"]}),
                (BLANK_STR,  {"size": 12, "color": C["BLACK"]}),
                (after,      {"size": 12, "color": C["BLACK"]}),
            ], 1.95, CONT_Y + 0.05 + i * row_h, 7.75, row_h,
               align="left", valign="middle", margin=0.03)

    def slide17():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Cloze", "Answers", f"{CODE}.16")

        row_h = (CONT_H - 0.10) / len(WORDS)
        wb_rows = [[{"text": w, "fill": C["WHITE"], "align": "center",
                     "size": 12, "valign": "middle"}]
                   for w in WORDS]
        table(s, wb_rows,
              0.30, CONT_Y + 0.05, 1.50, CONT_H - 0.10,
              [1.50], row_h)

        for i, w in enumerate(CLOZE_ORDER):
            sentence = SENTENCES[w]
            idx = sentence.find(w)
            before = sentence[:idx] if idx != -1 else sentence
            after  = sentence[idx + len(w):] if idx != -1 else ""
            rich_txt(s, [
                (before, {"size": 12, "color": C["BLACK"]}),
                (w,      {"size": 12, "bold": True, "color": C["PINK"]}),
                (after,  {"size": 12, "color": C["BLACK"]}),
            ], 1.95, CONT_Y + 0.05 + i * row_h, 7.75, row_h,
               align="left", valign="middle", margin=0.03)

    def _word_shed_structure(s, show_answers=False):
        slide_num = f"{CODE}.{'18' if show_answers else '17'}"
        add_frame(s, "Independent", "Word Shed",
                  "Build everything you can around the base word.", slide_num)

        FY = CONT_Y + 0.05
        FH = CONT_H - 0.10
        rect(s, 0.6, FY, 8.8, FH, "C4956A", "8B5E3C", 2)
        rect(s, 0.85, FY + 0.18, 8.3, FH - 0.36, C["WHITE"], "8B5E3C", 1)

        cx_split = 5.0
        body_x_l = 0.85
        body_x_r = cx_split + 0.05
        body_w = (cx_split - 0.05) - body_x_l
        body_y_top = FY + 0.18
        body_y_bot = FY + FH - 0.18
        body_h = body_y_bot - body_y_top
        mid_y = body_y_top + body_h / 2

        rect(s, cx_split - 0.02, body_y_top, 0.04, body_h, "AAAAAA")
        rect(s, body_x_l, mid_y - 0.02, (body_x_r + body_w) - body_x_l, 0.04, "AAAAAA")

        rect(s, cx_split - 1.1, mid_y - 0.28, 2.2, 0.56,
             C["WHITE"], C["BLACK"], 1.5)
        txt(s, lesson["wordShed"]["baseWord"],
            cx_split - 1.1, mid_y - 0.28, 2.2, 0.56,
            size=20, bold=True, color=C["BLACK"], align="center", margin=0)

        ws = lesson["wordShed"]
        bottom_label_y = mid_y + 0.32
        quadrants = [
            ("Definition",                body_x_l, body_y_top + 0.05, body_w,
             ws["def"] if show_answers else ""),
            ("In a Sentence",             body_x_r, body_y_top + 0.05, body_w,
             ws["sentence"] if show_answers else ""),
            ("Rhymes With...",            body_x_l, bottom_label_y, body_w,
             ws["rhymes"] if show_answers else ""),
            ("Add Prefixes or Suffixes",  body_x_r, bottom_label_y, body_w,
             ws["morphology"] if show_answers else ""),
        ]
        for label, lx, ly, lw, body_text in quadrants:
            txt(s, label, lx, ly, lw, 0.30,
                size=13, bold=True, color=C["BLACK"], align="left", valign="top", margin=0)
            if body_text:
                bx, bw = lx, lw
                if ly < mid_y:
                    by = ly + 0.34
                    bh = (mid_y - 0.32) - by
                else:
                    by = ly + 0.34
                    bh = (body_y_bot - 0.05) - by
                txt(s, body_text, bx, by, bw, bh,
                    size=12, color=C["PINK"], align="left", valign="top")

    def slide18():
        s = prs.slides.add_slide(BLANK)
        _word_shed_structure(s, show_answers=False)

    def slide19():
        s = prs.slides.add_slide(BLANK)
        _word_shed_structure(s, show_answers=True)

    def _morph_matrix(s, show_answers=False, sn_blank="19", sn_answers="20"):
        add_frame(s, "Whole Group", "Morphology Matrix",
                  "How many new words can you create by adding suffix(es)?",
                  f"{CODE}.{sn_answers if show_answers else sn_blank}")

        mm = lesson["morphMatrix"]

        txt(s, "Base Word", 0.4, CONT_Y + 0.15, 4.2, 0.38,
            size=20, bold=True, color=C["GREEN_S"], align="center")
        rect(s, 0.4, CONT_Y + 0.55, 4.2, 3.45, C["WHITE"], C["GREEN_S"], 2.5)
        txt(s, mm["baseWord"], 0.4, CONT_Y + 0.85, 4.2, 1.5,
            size=64, color=C["BLACK"], align="center", valign="middle")
        txt(s, mm["def"], 0.6, CONT_Y + 2.55, 3.8, 0.6,
            size=13, italic=True, color=C["GREY"], align="center", valign="middle")

        txt(s, "Suffix", 5.0, CONT_Y + 0.15, 4.6, 0.38,
            size=20, bold=True, color=C["PURPLE"], align="center")

        gcw, grh = 2.3, 1.10
        for i, sf in enumerate(mm["suffixes"]):
            col, row = i % 2, i // 2
            x = 5.0 + col * gcw
            y = CONT_Y + 0.55 + row * grh
            rect(s, x, y, gcw, grh, C["WHITE"], C["PURPLE"], 2)
            if show_answers:
                ans = mm["answers"][i] if i < len(mm["answers"]) else ""
                txt(s, sf, x, y + 0.05, gcw, 0.28,
                    size=12, color=C["GREY"], align="center", valign="top", margin=0)
                txt(s, ans, x + 0.06, y + 0.34, gcw - 0.12, grh - 0.40,
                    size=14, bold=True, color=C["PINK"],
                    align="center", valign="middle", margin=0)
            else:
                txt(s, sf, x, y, gcw, grh,
                    size=24, color=C["BLACK"], align="center", valign="middle", margin=0)

    def slide20():
        s = prs.slides.add_slide(BLANK)
        _morph_matrix(s, show_answers=False, sn_blank="19", sn_answers="20")

    def slide21():
        s = prs.slides.add_slide(BLANK)
        _morph_matrix(s, show_answers=True, sn_blank="19", sn_answers="20")

    # ════════════════════════════════════════════════════════════════════════
    # Y2 SLIDES
    # ════════════════════════════════════════════════════════════════════════

    def slide_y2_starter():
        """Y2 Starter: written word clues, click-reveal answers."""
        starter = lesson.get("y2Starter", {})
        clues = starter.get("clues", [])
        question = starter.get("question", "Which of last week's words is being described?")

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Starter", question, f"{CODE}.1")

        cell_w = 9.0 / 3
        cell_h = 1.55
        sx, sy = 0.5, CONT_Y + 0.18

        click_groups = []
        for i, item in enumerate(clues[:6]):
            col, row = i % 3, i // 3
            x = sx + col * cell_w
            y = sy + row * cell_h
            # Clue box
            rect(s, x + 0.06, y, cell_w - 0.12, 0.88, "F0F4FF", C["BLUE"], 0.8, radius=0.08)
            txt(s, item.get("clue", ""), x + 0.12, y + 0.04, cell_w - 0.24, 0.80,
                size=13, color=C["BLACK"], align="center", valign="middle")
            # Answer (hidden, click-reveal)
            ans_shape = txt(s, item.get("word", ""), x, y + 0.95, cell_w, 0.48,
                size=28, bold=True, color=C["PINK"], align="center")
            click_groups.append([ans_shape])

        register_clicks(s, *click_groups)

    def slide_y2_spelling_pattern():
        """Y2 Spelling Pattern: rule explanation with examples."""
        sp = lesson.get("spellingPattern", {})
        title = sp.get("title", lesson.get("rule", "Spelling Pattern"))
        body  = sp.get("body", "")
        examples = sp.get("examples", [])[:3]

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Spelling Pattern", title, f"{CODE}.3")

        txt(s, body, 0.5, CONT_Y + 0.18, 9.0, 0.75,
            size=17, color=C["BLACK"], align="center", valign="middle")

        if examples:
            n = len(examples)
            bw = (9.0 - 0.3 * (n - 1)) / n
            ex_y = CONT_Y + 1.05
            ex_h = 2.85

            click_shapes = []
            for i, ex in enumerate(examples):
                x = 0.5 + i * (bw + 0.3)
                rect(s, x, ex_y, bw, ex_h, C["WHITE"], C["BLUE"], 1.5, radius=0.1)

                inp = ex.get("input", ex.get("base", ""))
                out = ex.get("output", ex.get("result", ""))
                note = ex.get("note", "")

                if inp:
                    txt(s, inp, x + 0.08, ex_y + 0.12, bw - 0.16, 0.55,
                        size=26, color=C["BLACK"], align="center")
                    txt(s, "↓", x, ex_y + 0.72, bw, 0.35,
                        size=22, color=C["GREY"], align="center")

                out_shape = txt(s, out, x + 0.08, ex_y + 1.12, bw - 0.16, 0.65,
                    size=28, bold=True, color=C["PINK"], align="center")
                click_shapes.append(out_shape)

                if note:
                    note_shape = txt(s, note, x + 0.1, ex_y + 1.85, bw - 0.2, 0.85,
                        size=12, italic=True, color=C["GREY"], align="center", valign="top")
                    click_shapes.append(note_shape)

            register_clicks(s, click_shapes)

    def slide_y2_morph_matrix_blank():
        s = prs.slides.add_slide(BLANK)
        _morph_matrix(s, show_answers=False, sn_blank="4a", sn_answers="4b")

    def slide_y2_morph_matrix_answers():
        s = prs.slides.add_slide(BLANK)
        _morph_matrix(s, show_answers=True, sn_blank="4a", sn_answers="4b")

    def slide_y2_word_sort_blank():
        """Y2 Word Sort (blank) — same as Y4 slide7 but with Y2 slide number."""
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Word Sort", lesson["wordSortQ"], f"{CODE}.4")

        for ri, row_words in enumerate([WORDS[:5], WORDS[5:]]):
            txt(s, "          ".join(row_words),
                0.5, CONT_Y + 0.15 + ri * 0.5, 9.0, 0.48,
                size=22, color=C["BLACK"], align="center")

        BOX_Y, BOX_H = CONT_Y + 1.3, 2.05
        ws = lesson["wordSort"]
        rect(s, 0.35, BOX_Y, 4.45, BOX_H, C["WHITE"], C["PURPLE"], 3, radius=0.15)
        txt(s, ws["box1label"] + "\n" + ws["box1sub"],
            0.45, BOX_Y + 0.08, 4.25, 0.75,
            size=16, bold=True, color=C["PURPLE"], align="center")
        rect(s, 5.2, BOX_Y, 4.45, BOX_H, C["WHITE"], C["GREEN_S"], 3, radius=0.15)
        txt(s, ws["box2label"] + "\n" + ws["box2sub"],
            5.3, BOX_Y + 0.08, 4.25, 0.75,
            size=16, bold=True, color=C["GREEN_S"], align="center")
        txt(s, ws["hint"], 0.5, BOX_Y + BOX_H + 0.12, 9.0, 0.32,
            size=12, italic=True, color=C["GREY"], align="center")

    def slide_y2_word_sort_answers():
        """Y2 Word Sort answers — same as Y4 slide8 but with Y2 slide number."""
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Word Sort", "Answers", f"{CODE}.5")

        ws = lesson["wordSort"]
        txt(s, ws["answerNote"], 0.5, CONT_Y + 0.08, 9.0, 0.38,
            size=14, color=C["BLACK"], align="center")

        BOX_Y, BOX_H = CONT_Y + 0.55, 3.05
        verb_only = ws["verbOnly"]
        verb_noun = ws["verbNoun"]

        rect(s, 0.35, BOX_Y, 4.45, BOX_H, C["WHITE"], C["PURPLE"], 3, radius=0.15)
        txt(s, ws["box1label"] + "  " + ws["box1sub"],
            0.45, BOX_Y + 0.08, 4.25, 0.35,
            size=14, bold=True, color=C["PURPLE"], align="center")
        if verb_only:
            available_h = BOX_H - 0.50 - 0.10
            row_h  = min(0.52, available_h / len(verb_only))
            font_s = max(9, min(13, int(row_h * 22)))
            for i, w in enumerate(verb_only):
                txt(s, w, 0.55, BOX_Y + 0.50 + i * row_h, 4.15, row_h,
                    size=font_s, bold=True, color=C["PINK"], align="center")

        rect(s, 5.2, BOX_Y, 4.45, BOX_H, C["WHITE"], C["GREEN_S"], 3, radius=0.15)
        txt(s, ws["box2label"], 5.3, BOX_Y + 0.08, 4.25, 0.35,
            size=14, bold=True, color=C["GREEN_S"], align="center")
        if verb_noun:
            available_h = BOX_H - 0.47 - 0.10
            row_h  = min(0.362, available_h / len(verb_noun))
            font_w = max(9, min(12, int(row_h * 28)))
            font_eg = max(8, font_w - 2)
            for i, item in enumerate(verb_noun):
                rich_txt(s, [
                    (item["word"], {"size": font_w, "bold": True, "color": C["PINK"]}),
                    ("  " + item["eg"], {"size": font_eg, "italic": True, "color": C["GREY"]}),
                ], 5.4, BOX_Y + 0.47 + i * row_h, 4.15, row_h, valign="middle")

        txt(s, ws["exampleLine"], 0.5, BOX_Y + BOX_H + 0.12, 9.0, 0.30,
            size=12, color=C["BLACK"], align="center")

    def slide_y2_sound_it_squash_it():
        """Y2 Sound It, Squash It, Say It, Scribe It — worked examples with sound buttons."""
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Sound It, Squash It, Say It, Scribe It",
                  "Can you sound out and write the words? Can you add the sound buttons underneath?",
                  f"{CODE}.6")

        txt(s, "Use a dot ● for each single letter sound and a line — when two or more letters make one sound.",
            0.5, CONT_Y + 0.12, 9.0, 0.30,
            size=12, italic=True, color=C["GREY"], align="center")

        example_words = WORDS[:4]
        n_ex = len(example_words)
        cw = 9.0 / n_ex
        sy = CONT_Y + 0.52
        syllables = lesson.get("wordMaps", {}).get("syllables", {})
        fs = min(fit_font(w, cw, 32, 20) for w in example_words)

        for i, w in enumerate(example_words):
            x  = 0.5 + i * cw
            cx = x + cw / 2
            txt(s, w, x, sy, cw, 0.72,
                size=fs, color=C["PINK"], align="center")
            draw_sound_buttons(s, w, cx, sy + 0.80, 22)
            brk = syllables.get(w, "")
            if brk:
                txt(s, brk, x, sy + 1.68, cw, 0.28,
                    size=12, italic=True, color=C["GREY"], align="center")

        # Practice boxes for next 2 words
        practice_words = WORDS[4:6] if len(WORDS) > 4 else []
        if practice_words:
            txt(s, "Now try these:", 0.5, CONT_Y + 2.62, 3.0, 0.28,
                size=13, bold=True, color=C["BLUE"], align="left")
            pw = 4.2
            for i, w in enumerate(practice_words):
                px = 0.5 + i * (pw + 0.15)
                rect(s, px, CONT_Y + 2.98, pw, 1.05, C["WHITE"], C["GREY"], 0.8, radius=0.06)
                txt(s, w, px, CONT_Y + 3.02, pw, 0.40,
                    size=24, color=C["BLACK"], align="center")
                txt(s, "(add sound buttons below)", px, CONT_Y + 3.46, pw, 0.28,
                    size=10, italic=True, color=C["GREY"], align="center")

    def slide_y2_phoneme_maps_blank():
        """Y2 Phoneme Maps — symbols only, word hidden."""
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Phoneme Maps",
                  "Look at the phoneme maps. Write the matching word next to each map.",
                  f"{CODE}.7")

        txt(s, "Each ● = one sound.  Each — = two or more letters making one sound.",
            0.5, CONT_Y + 0.10, 9.0, 0.25,
            size=11, italic=True, color=C["GREY"], align="center")

        # 5 words per column, 2 columns
        col_x = [0.55, 5.1]
        row_h = (CONT_H - 0.55) / 5
        sym_w = 3.2   # width of symbol area
        lbl_w = 1.5   # width of write-word area
        scale = 1.0

        for i, w in enumerate(WORDS[:10]):
            col = i // 5
            row = i % 5
            x   = col_x[col]
            y   = CONT_Y + 0.42 + row * row_h
            cx  = x + sym_w / 2
            sym_y = y + (row_h - 0.07) / 2

            draw_phoneme_symbols_only(s, w, cx, sym_y, scale=scale)

            # Writing line for the word
            rect(s, x + sym_w + 0.12, y + row_h * 0.28, lbl_w, 0.03, C["GREY"])
            txt(s, "_____________", x + sym_w + 0.12, y + row_h * 0.05, lbl_w, row_h * 0.6,
                size=13, color="CCCCCC", align="left", valign="bottom")

    def slide_y2_phoneme_maps_answers():
        """Y2 Phoneme Maps answers — word shown next to each symbol set."""
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Phoneme Maps", "Answers", f"{CODE}.8")

        txt(s, "Each ● = one sound.  Each — = two or more letters making one sound.",
            0.5, CONT_Y + 0.10, 9.0, 0.25,
            size=11, italic=True, color=C["GREY"], align="center")

        col_x = [0.55, 5.1]
        row_h = (CONT_H - 0.55) / 5
        sym_w = 3.2
        lbl_w = 1.5
        scale = 1.0

        click_shapes = []
        for i, w in enumerate(WORDS[:10]):
            col = i // 5
            row = i % 5
            x   = col_x[col]
            y   = CONT_Y + 0.42 + row * row_h
            cx  = x + sym_w / 2
            sym_y = y + (row_h - 0.07) / 2

            draw_phoneme_symbols_only(s, w, cx, sym_y, scale=scale)

            ans_shape = txt(s, w, x + sym_w + 0.12, y + row_h * 0.08,
                lbl_w, row_h * 0.82,
                size=16, bold=True, color=C["PINK"], align="left", valign="middle")
            click_shapes.append(ans_shape)

        register_clicks(s, click_shapes)

    def slide_y2_sentences_synonyms_blank():
        """Y2 Sentences and Synonyms — synonym replacement task (blank)."""
        sns = lesson.get("sentencesAndSynonyms", [])[:4]

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Sentences and Synonyms",
                  "Read each sentence. Replace the underlined word with one from the word bank.",
                  f"{CODE}.9")

        # Word bank at top in blue
        wb_y = CONT_Y + 0.12
        txt(s, "Word bank:", 0.5, wb_y, 1.4, 0.32,
            size=13, bold=True, color=C["BLUE"], align="left")
        txt(s, "   ".join(WORDS), 1.9, wb_y, 7.6, 0.32,
            size=13, bold=True, color=C["BLUE"], align="left")

        # Sentence pairs in 2x2 grid
        grid_y = CONT_Y + 0.55
        cw, rh = 4.55, 1.72
        for i, item in enumerate(sns):
            col, row = i % 2, i // 2
            x = 0.35 + col * (cw + 0.25)
            y = grid_y + row * (rh + 0.10)
            rect(s, x, y, cw, rh, C["WHITE"], "CCCCCC", 0.8, radius=0.06)

            original = item.get("original", "")
            underlined = item.get("underlined", "")

            # Split sentence around underlined word
            idx = original.find(underlined)
            if idx != -1:
                before = original[:idx]
                after  = original[idx + len(underlined):]
                rich_txt(s, [
                    (before,      {"size": 12, "color": C["BLACK"]}),
                    (underlined,  {"size": 12, "color": C["BLACK"], "bold": True}),
                    (after,       {"size": 12, "color": C["BLACK"]}),
                ], x + 0.12, y + 0.08, cw - 0.24, 0.65,
                   align="left", valign="top", margin=0.02)
            else:
                txt(s, original, x + 0.12, y + 0.08, cw - 0.24, 0.65,
                    size=12, color=C["BLACK"], align="left", valign="top")

            txt(s, f"Replace the word in bold with: ______________",
                x + 0.12, y + 0.80, cw - 0.24, 0.38,
                size=11, color=C["GREY"], align="left")

    def slide_y2_sentences_synonyms_answers():
        """Y2 Sentences and Synonyms answers."""
        sns = lesson.get("sentencesAndSynonyms", [])[:4]

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Sentences and Synonyms", "Answers", f"{CODE}.10")

        txt(s, "Word bank:", 0.5, CONT_Y + 0.12, 1.4, 0.32,
            size=13, bold=True, color=C["BLUE"], align="left")
        txt(s, "   ".join(WORDS), 1.9, CONT_Y + 0.12, 7.6, 0.32,
            size=13, bold=True, color=C["BLUE"], align="left")

        grid_y = CONT_Y + 0.55
        cw, rh = 4.55, 1.72
        click_shapes = []
        for i, item in enumerate(sns):
            col, row = i % 2, i // 2
            x = 0.35 + col * (cw + 0.25)
            y = grid_y + row * (rh + 0.10)
            rect(s, x, y, cw, rh, C["WHITE"], "CCCCCC", 0.8, radius=0.06)

            original = item.get("original", "")
            underlined = item.get("underlined", "")
            answer = item.get("answer", "")

            idx = original.find(underlined)
            if idx != -1:
                before = original[:idx]
                after  = original[idx + len(underlined):]
                rich_txt(s, [
                    (before,     {"size": 12, "color": C["BLACK"]}),
                    (underlined, {"size": 12, "color": C["BLACK"], "bold": True}),
                    (after,      {"size": 12, "color": C["BLACK"]}),
                ], x + 0.12, y + 0.08, cw - 0.24, 0.65,
                   align="left", valign="top", margin=0.02)
            else:
                txt(s, original, x + 0.12, y + 0.08, cw - 0.24, 0.65,
                    size=12, color=C["BLACK"], align="left", valign="top")

            ans_shape = txt(s, f"→  {answer}", x + 0.12, y + 0.82, cw - 0.24, 0.52,
                size=18, bold=True, color=C["PINK"], align="left", valign="middle")
            click_shapes.append(ans_shape)

        register_clicks(s, click_shapes)

    def slide_y2_words_in_action_blank():
        """Y2 Words in Action — creative writing prompt."""
        wia = lesson.get("wordsInAction", {})
        prompt = wia.get("prompt", "Can you write a description of this picture?")
        pic_prompt = wia.get("picture_prompt", "")
        req_words = wia.get("required_words", WORDS[:3])

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Words in Action", prompt, f"{CODE}.11")

        # Picture frame (left side)
        rect(s, 0.35, CONT_Y + 0.18, 5.2, 3.55, "F9F9F9", C["GREY"], 0.8)
        txt(s, "🖼", 0.35, CONT_Y + 0.22, 5.2, 0.55,
            size=28, color=C["GREY"], align="center")
        txt(s, pic_prompt, 0.5, CONT_Y + 0.82, 4.9, 2.85,
            size=13, italic=True, color=C["GREY"], align="center", valign="top")

        # Required words (right side)
        txt(s, "Words to include:", 5.75, CONT_Y + 0.22, 3.9, 0.35,
            size=14, bold=True, color=C["BLUE"], align="left")
        for i, w in enumerate(req_words):
            txt(s, w, 5.75, CONT_Y + 0.65 + i * 0.52, 3.9, 0.45,
                size=22, bold=True, color=C["PINK"], align="left")

    def slide_y2_words_in_action_answers():
        """Y2 Words in Action answers — possible answer shown."""
        wia = lesson.get("wordsInAction", {})
        prompt = wia.get("prompt", "Can you write a description of this picture?")
        req_words = wia.get("required_words", WORDS[:3])
        sample = wia.get("sample_answer", "")

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Words in Action",
                  "Possible answer — pupils' answers will vary!", f"{CODE}.12")

        if sample:
            txt(s, sample, 0.5, CONT_Y + 0.25, 9.0, 3.0,
                size=16, color=C["BLACK"], align="left", valign="top")

            # Highlight the required words
            txt(s, "Required words used:  " + "   ".join(req_words),
                0.5, CONT_Y + 3.45, 9.0, 0.35,
                size=13, bold=True, color=C["PINK"], align="center")
        else:
            txt(s, "Pupils write their own description using the required words.",
                0.5, CONT_Y + 1.5, 9.0, 0.5,
                size=16, italic=True, color=C["GREY"], align="center")
            txt(s, "Words to use:  " + "   ".join(req_words),
                0.5, CONT_Y + 2.2, 9.0, 0.45,
                size=18, bold=True, color=C["PINK"], align="center")

    def slide_y2_word_spotter():
        """Y2 Word Spotter — 4x4 grid, clap for this week's words."""
        distractors = lesson.get("wordSpotter", {}).get("distractors", [])

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Word Spotter",
                  "Clap your hands when you see one of this week's words.", f"{CODE}.13")

        all_words = list(WORDS) + list(distractors[:6])
        # Shuffle with fixed seed for reproducibility
        rng = random.Random(42)
        rng.shuffle(all_words)
        # Pad or trim to exactly 16
        while len(all_words) < 16:
            all_words.append(all_words[0])
        all_words = all_words[:16]

        this_week_set = set(WORDS)
        COLS, ROWS = 4, 4
        cell_w = 9.0 / COLS
        cell_h = (CONT_H - 0.2) / ROWS
        sx, sy = 0.5, CONT_Y + 0.1

        for i, w in enumerate(all_words):
            col, row = i % COLS, i // COLS
            x = sx + col * cell_w
            y = sy + row * cell_h
            is_this_week = w in this_week_set
            bg = "FFF0F5" if is_this_week else C["WHITE"]
            border = C["PINK"] if is_this_week else "CCCCCC"
            lw = 1.5 if is_this_week else 0.5
            rect(s, x + 0.05, y + 0.05, cell_w - 0.1, cell_h - 0.1, bg, border, lw, radius=0.08)
            txt(s, w, x + 0.05, y + 0.05, cell_w - 0.1, cell_h - 0.1,
                size=18, color=C["BLACK"], align="center", valign="middle")

    def slide_y2_continuous_provision():
        """Y2 Continuous Provision — word cards with phoneme symbols."""
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Optional", "Continuous Provision",
                  "Word cards for independent learning", f"{CODE}.14")

        COLS, ROWS = 5, 2
        cw = 9.2 / COLS
        ch = (CONT_H - 0.15) / ROWS
        sx, sy = 0.4, CONT_Y + 0.08

        for i, w in enumerate(WORDS[:10]):
            col, row = i % COLS, i // COLS
            x = sx + col * cw
            y = sy + row * ch
            rect(s, x + 0.04, y + 0.04, cw - 0.08, ch - 0.08,
                 C["WHITE"], "CCCCCC", 0.6, radius=0.08)
            word_h = ch * 0.48
            sym_h  = ch * 0.38
            txt(s, w, x + 0.06, y + 0.06, cw - 0.12, word_h,
                size=min(fit_font(w, cw, 24, 14), 22),
                color=C["BLACK"], align="center", valign="middle")
            cx = x + cw / 2
            draw_phoneme_symbols_only(s, w, cx, y + word_h + 0.08, scale=0.9)

    # ════════════════════════════════════════════════════════════════════════
    # Y3 SLIDES
    # ════════════════════════════════════════════════════════════════════════

    def slide_y3_starter_blank():
        """Y3 Starter: sorting mat (blank)."""
        starter = lesson.get("y3Starter", {})
        question = starter.get("question", "Can you sort the words based on their suffix?")
        categories = starter.get("categories", [])
        words_data = starter.get("last_week_words", [])

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Sorting Mat", question, f"{CODE}.1")

        if not categories or not words_data:
            txt(s, "Sort last week's words into the correct categories.",
                0.5, CONT_Y + 1.5, 9.0, 0.5,
                size=16, color=C["BLACK"], align="center")
            return

        n_cats = len(categories)
        cat_w  = (9.0 - 0.3) / n_cats
        hdr_y  = CONT_Y + 0.15
        hdr_h  = 0.45
        body_y = hdr_y + hdr_h + 0.08
        body_h = CONT_H - hdr_h - 0.32

        cat_colors = [C["PURPLE"], C["GREEN_S"], C["BLUE"], C["PINK"]]

        for ci, cat in enumerate(categories):
            cx = 0.5 + ci * (cat_w + 0.1)
            color = cat_colors[ci % len(cat_colors)]
            rect(s, cx, hdr_y, cat_w, hdr_h, color, color, 0)
            txt(s, cat, cx, hdr_y, cat_w, hdr_h,
                size=20, bold=True, color=C["WHITE"], align="center", margin=0)
            rect(s, cx, body_y, cat_w, body_h, C["WHITE"], color, 1.5, radius=0.08)

        # Word bank below header in a separate area
        all_wds = [item.get("word", "") for item in words_data]
        txt(s, "Words to sort: " + "  |  ".join(all_wds),
            0.5, BAR_Y - 0.42, 9.0, 0.30,
            size=12, color=C["BLACK"], align="center")

    def slide_y3_starter_answers():
        """Y3 Starter: sorting mat answers (click-reveal groups by category)."""
        starter = lesson.get("y3Starter", {})
        question = starter.get("question", "Can you sort the words based on their suffix?")
        categories = starter.get("categories", [])
        words_data = starter.get("last_week_words", [])

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Sorting Mat", "Answers", f"{CODE}.2")

        if not categories or not words_data:
            return

        n_cats = len(categories)
        cat_w  = (9.0 - 0.3) / n_cats
        hdr_y  = CONT_Y + 0.15
        hdr_h  = 0.45
        body_y = hdr_y + hdr_h + 0.08
        body_h = CONT_H - hdr_h - 0.32

        cat_colors = [C["PURPLE"], C["GREEN_S"], C["BLUE"], C["PINK"]]

        for ci, cat in enumerate(categories):
            cx = 0.5 + ci * (cat_w + 0.1)
            color = cat_colors[ci % len(cat_colors)]
            rect(s, cx, hdr_y, cat_w, hdr_h, color, color, 0)
            txt(s, cat, cx, hdr_y, cat_w, hdr_h,
                size=20, bold=True, color=C["WHITE"], align="center", margin=0)
            rect(s, cx, body_y, cat_w, body_h, C["WHITE"], color, 1.5, radius=0.08)

            cat_words = [item.get("word", "") for item in words_data
                        if item.get("category", "") == cat]
            n = len(cat_words)
            if n:
                row_h = min(0.48, (body_h - 0.12) / n)
                click_batch = []
                for ri, w in enumerate(cat_words):
                    ws = txt(s, w, cx + 0.08, body_y + 0.08 + ri * row_h,
                        cat_w - 0.16, row_h,
                        size=max(12, min(18, int(row_h * 28))),
                        bold=True, color=C["PINK"], align="center", valign="middle")
                    click_batch.append(ws)
                register_clicks(s, click_batch)

    def slide_y3_this_weeks_words():
        """Y3 This Week's Words — same layout as Y4 but Y3 slide number."""
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "This Week's Words",
                  lesson["thisWeeksWordsQ"], f"{CODE}.3")

        cell_w = 9.0 / 5
        fs = min(fit_font(w, cell_w, 32, 18) for w in WORDS)
        for ri, row_words in enumerate([WORDS[:5], WORDS[5:]]):
            for ci, w in enumerate(row_words):
                x = 0.5 + ci * cell_w
                y = CONT_Y + 0.3 + ri * 1.30
                txt(s, w, x, y, cell_w, 0.85,
                    size=fs, color=C["BLACK"], align="center", valign="bottom")
                rect(s, x + 0.3, y + 0.92, cell_w - 0.6, 0.06, C["PINK"])

        prompt_shape = txt(s, lesson["thisWeeksWordsPrompt"],
            0.5, CONT_Y + 3.05, 4.3, 0.4,
            size=15, bold=True, color=C["BLUE"], align="left")
        explanation_shape = txt(s, lesson["thisWeeksWordsExplanation"],
            4.8, CONT_Y + 3.05, 4.8, 0.55,
            size=14, color=C["BLACK"], align="left")
        register_clicks(s, prompt_shape, explanation_shape)

    def slide_y3_spelling_pattern():
        """Y3 Spelling Pattern — rule explanation with transformation examples."""
        sp = lesson.get("spellingPattern", {})
        title    = sp.get("title", lesson.get("rule", "Spelling Pattern"))
        body     = sp.get("body", "")
        rule_note = sp.get("rule_note", "")
        examples = sp.get("examples", [])[:3]

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Spelling Pattern", title, f"{CODE}.4")

        txt(s, body, 0.5, CONT_Y + 0.14, 9.0, 0.62,
            size=16, color=C["BLACK"], align="center", valign="middle")

        if examples:
            n = len(examples)
            bw = (9.0 - 0.25 * (n - 1)) / n
            ex_y = CONT_Y + 0.90
            ex_h = 2.70

            for i, ex in enumerate(examples):
                x = 0.5 + i * (bw + 0.25)
                rect(s, x, ex_y, bw, ex_h, C["WHITE"], C["BLUE"], 1.5, radius=0.10)

                base   = ex.get("base", ex.get("input", ""))
                result = ex.get("result", ex.get("output", ""))
                note   = ex.get("note", "")

                if base:
                    txt(s, base, x + 0.1, ex_y + 0.12, bw - 0.2, 0.50,
                        size=24, color=C["BLACK"], align="center")
                    txt(s, "↓", x, ex_y + 0.66, bw, 0.32,
                        size=20, color=C["GREY"], align="center")
                txt(s, result, x + 0.1, ex_y + 1.02, bw - 0.2, 0.58,
                    size=26, bold=True, color=C["PINK"], align="center")
                if note:
                    txt(s, note, x + 0.12, ex_y + 1.68, bw - 0.24, 0.90,
                        size=12, italic=True, color=C["GREY"], align="center", valign="top")

        if rule_note:
            txt(s, rule_note, 0.5, BAR_Y - 0.42, 9.0, 0.30,
                size=12, italic=True, color=C["GREY"], align="center")

    def slide_y3_etymology_simple():
        """Y3 Etymology — simplified single-slide version, not a guessing game."""
        etym = lesson.get("etymology", {})
        clicks = etym.get("clicks", [])
        base   = etym.get("baseForm", etym.get("word", ""))

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Etymology",
                  f"Where does the word '{base}' come from?", f"{CODE}.5")

        # Show the word prominently
        txt(s, base, 0.5, CONT_Y + 0.08, 9.0, 0.72,
            size=52, bold=True, color=C["PINK"], align="center")

        # Show up to 3 etymology fact panels
        panels = clicks[:3] if len(clicks) >= 3 else clicks
        n = len(panels)
        if n:
            pw = (9.0 - 0.15 * (n - 1)) / n
            py = CONT_Y + 0.90
            ph = 2.55

            for i, card in enumerate(panels):
                x = 0.5 + i * (pw + 0.15)
                rect(s, x, py, pw, ph, C["WHITE"], C["BLUE"], 1.5, radius=0.10)
                lbl  = card.get("label", "")
                body = card.get("body", "")
                if lbl:
                    txt(s, lbl, x + 0.08, py + 0.10, pw - 0.16, 0.38,
                        size=13, bold=True, color=C["BLUE"], align="center")
                txt(s, body, x + 0.10, py + 0.52, pw - 0.20, ph - 0.62,
                    size=12, color=C["BLACK"], align="center", valign="top")

        # Fun fact (click 4 or 5 if present)
        if len(clicks) > 3:
            extra = clicks[3]
            txt(s, extra.get("label", "") + ": " + extra.get("body", ""),
                0.5, BAR_Y - 0.42, 9.0, 0.30,
                size=11, italic=True, color=C["GREY"], align="center")

    def slide_y3_morph_matrix_blank():
        s = prs.slides.add_slide(BLANK)
        _morph_matrix(s, show_answers=False, sn_blank="6", sn_answers="7")

    def slide_y3_morph_matrix_answers():
        s = prs.slides.add_slide(BLANK)
        _morph_matrix(s, show_answers=True, sn_blank="6", sn_answers="7")

    def slide_y3_word_sort_blank():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Word Sort", lesson["wordSortQ"], f"{CODE}.8")
        for ri, row_words in enumerate([WORDS[:5], WORDS[5:]]):
            txt(s, "          ".join(row_words),
                0.5, CONT_Y + 0.15 + ri * 0.5, 9.0, 0.48,
                size=22, color=C["BLACK"], align="center")
        BOX_Y, BOX_H = CONT_Y + 1.3, 2.05
        ws = lesson["wordSort"]
        rect(s, 0.35, BOX_Y, 4.45, BOX_H, C["WHITE"], C["PURPLE"], 3, radius=0.15)
        txt(s, ws["box1label"] + "\n" + ws["box1sub"],
            0.45, BOX_Y + 0.08, 4.25, 0.75,
            size=16, bold=True, color=C["PURPLE"], align="center")
        rect(s, 5.2, BOX_Y, 4.45, BOX_H, C["WHITE"], C["GREEN_S"], 3, radius=0.15)
        txt(s, ws["box2label"] + "\n" + ws["box2sub"],
            5.3, BOX_Y + 0.08, 4.25, 0.75,
            size=16, bold=True, color=C["GREEN_S"], align="center")
        txt(s, ws["hint"], 0.5, BOX_Y + BOX_H + 0.12, 9.0, 0.32,
            size=12, italic=True, color=C["GREY"], align="center")

    def slide_y3_word_sort_answers():
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Word Sort", "Answers", f"{CODE}.9")
        ws = lesson["wordSort"]
        txt(s, ws["answerNote"], 0.5, CONT_Y + 0.08, 9.0, 0.38,
            size=14, color=C["BLACK"], align="center")
        BOX_Y, BOX_H = CONT_Y + 0.55, 3.05
        verb_only = ws["verbOnly"]
        verb_noun = ws["verbNoun"]
        rect(s, 0.35, BOX_Y, 4.45, BOX_H, C["WHITE"], C["PURPLE"], 3, radius=0.15)
        txt(s, ws["box1label"] + "  " + ws["box1sub"],
            0.45, BOX_Y + 0.08, 4.25, 0.35,
            size=14, bold=True, color=C["PURPLE"], align="center")
        if verb_only:
            available_h = BOX_H - 0.50 - 0.10
            row_h  = min(0.52, available_h / len(verb_only))
            font_s = max(9, min(13, int(row_h * 22)))
            for i, w in enumerate(verb_only):
                txt(s, w, 0.55, BOX_Y + 0.50 + i * row_h, 4.15, row_h,
                    size=font_s, bold=True, color=C["PINK"], align="center")
        rect(s, 5.2, BOX_Y, 4.45, BOX_H, C["WHITE"], C["GREEN_S"], 3, radius=0.15)
        txt(s, ws["box2label"], 5.3, BOX_Y + 0.08, 4.25, 0.35,
            size=14, bold=True, color=C["GREEN_S"], align="center")
        if verb_noun:
            available_h = BOX_H - 0.47 - 0.10
            row_h  = min(0.362, available_h / len(verb_noun))
            font_w = max(9, min(12, int(row_h * 28)))
            font_eg = max(8, font_w - 2)
            for i, item in enumerate(verb_noun):
                rich_txt(s, [
                    (item["word"], {"size": font_w, "bold": True, "color": C["PINK"]}),
                    ("  " + item["eg"], {"size": font_eg, "italic": True, "color": C["GREY"]}),
                ], 5.4, BOX_Y + 0.47 + i * row_h, 4.15, row_h, valign="middle")
        txt(s, ws["exampleLine"], 0.5, BOX_Y + BOX_H + 0.12, 9.0, 0.30,
            size=12, color=C["BLACK"], align="center")

    def slide_y3_sound_buttons_syllables():
        """Y3 Sound Buttons and Syllables — worked examples, same as Y4 slide9 but Y3 number."""
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Sound Buttons and Syllables",
                  "How to mark the sound buttons and syllable breaks", f"{CODE}.10")

        map_words = lesson["wordMaps"]["words"][:3]
        syllables = lesson["wordMaps"]["syllables"]
        sc        = lesson["syllableCounts"]
        cw, sx, sy = 2.9, 0.5, CONT_Y + 0.25
        fs = min(fit_font(w, cw, 36, 18) for w in map_words)

        for i, w in enumerate(map_words):
            x  = sx + i * (cw + 0.25)
            cx = x + cw / 2
            brk = syllables.get(w, w)
            n   = sc.get(w, 1)
            txt(s, w, x, sy, cw, 0.7, size=fs, color=C["PINK"], align="center")
            txt(s, "Syllable Breaks", x, sy + 0.75, cw, 0.3,
                size=12, bold=True, color=C["BLUE"], align="center")
            rect(s, x, sy + 1.08, cw, 0.5, C["WHITE"], C["BLUE"], 1.5)
            txt(s, brk, x, sy + 1.08, cw, 0.5, size=20, color=C["BLACK"],
                align="center", margin=0)
            txt(s, f"{n} syllable{'s' if n > 1 else ''}",
                x, sy + 1.62, cw, 0.28, size=12, italic=True, color=C["GREY"], align="center")
            txt(s, "Sound Buttons", x, sy + 1.95, cw, 0.3,
                size=12, bold=True, color=C["PINK"], align="center")
            draw_sound_buttons(s, w, cx, sy + 2.32, 22)
            txt(s, "● = single sound   — = letters making one sound",
                x, sy + 2.88, cw, 0.22, size=9, italic=True, color=C["GREY"], align="center")

    def slide_y3_word_maps_blank():
        """Y3 Word Maps blank — same 4-column layout as Y4, Y3 slide number."""
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Word Maps",
                  "Read the word. Write syllable breaks and sound buttons. Rewrite the word.",
                  f"{CODE}.11")
        map_words = lesson["wordMaps"]["words"]
        hdr = [
            {"text": "Words",           "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "Syllable Breaks", "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "Sound Buttons",   "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "My Word",         "bold": True, "fill": C["WHITE"], "align": "center"},
        ]
        rows = [
            [{"text": w, "fill": C["WHITE"], "align": "center", "size": 18},
             {"text": "", "fill": C["WHITE"]},
             {"text": "", "fill": C["WHITE"]},
             {"text": "", "fill": C["WHITE"]}]
            for w in map_words
        ]
        rh = [0.5] + [0.57] * len(rows)
        table(s, [hdr] + rows,
              0.35, CONT_Y + 0.1, 9.3, CONT_H - 0.15,
              [2.0, 2.5, 2.9, 1.9], rh)

    def slide_y3_word_maps_answers():
        """Y3 Word Maps answers."""
        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Word Maps", "Answers", f"{CODE}.12")
        map_words = lesson["wordMaps"]["words"]
        syllables = lesson["wordMaps"]["syllables"]
        hdr = [
            {"text": "Words",           "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "Syllable Breaks", "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "Sound Buttons",   "bold": True, "fill": C["WHITE"], "align": "center"},
            {"text": "My Word",         "bold": True, "fill": C["WHITE"], "align": "center"},
        ]
        rows = [
            [{"text": w,                   "fill": C["WHITE"], "align": "center", "size": 18},
             {"text": syllables.get(w, w), "fill": C["WHITE"], "align": "center", "size": 16},
             {"text": "",                  "fill": C["WHITE"]},
             {"text": w,                   "fill": C["WHITE"], "align": "center", "size": 18, "color": C["PINK"]}]
            for w in map_words
        ]
        rh = [0.5] + [0.57] * len(rows)
        table(s, [hdr] + rows,
              0.35, CONT_Y + 0.1, 9.3, CONT_H - 0.15,
              [2.0, 2.5, 2.9, 1.9], rh)

        TABLE_X  = 0.35
        SB_COL_X = TABLE_X + 2.0 + 2.5
        SB_COL_W = 2.9
        HDR_H_IN = 0.5
        ROW_H_IN = 0.57
        SB_FS    = 13
        for ri, w in enumerate(map_words):
            cell_top_y = CONT_Y + 0.1 + HDR_H_IN + ri * ROW_H_IN
            cell_cx    = SB_COL_X + SB_COL_W / 2
            row_h_sb   = SB_FS * 1.4 / 72
            dot_h      = max(0.040, SB_FS * 0.084 / 28)
            total_sb_h = row_h_sb + 0.04 + dot_h
            word_top_y = cell_top_y + (ROW_H_IN - total_sb_h) / 2
            draw_sound_buttons(s, w, cell_cx, word_top_y, SB_FS)

    def slide_y3_word_match_blank():
        """Y3 Word Match — write the word that matches the description (blank)."""
        word_match = lesson.get("wordMatch", [])[:6]

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Word Match",
                  "Write the word that matches each description.", f"{CODE}.13")

        if not word_match:
            return

        n = len(word_match)
        row_h = (CONT_H - 0.15) / n
        for i, item in enumerate(word_match):
            y = CONT_Y + 0.08 + i * row_h
            desc = item.get("description", "")
            rect(s, 0.35, y + 0.04, 6.8, row_h - 0.08, C["WHITE"], "CCCCCC", 0.6)
            txt(s, desc, 0.5, y + 0.06, 6.6, row_h - 0.12,
                size=13, color=C["BLACK"], align="left", valign="middle")
            # Write-in box on right
            rect(s, 7.3, y + 0.12, 2.4, row_h - 0.24, "F9F9F9", C["GREY"], 0.8)
            txt(s, "____________", 7.35, y + 0.14, 2.3, row_h - 0.28,
                size=13, color="BBBBBB", align="center", valign="middle")

    def slide_y3_word_match_answers():
        """Y3 Word Match answers."""
        word_match = lesson.get("wordMatch", [])[:6]

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Word Match", "Answers", f"{CODE}.14")

        if not word_match:
            return

        n = len(word_match)
        row_h = (CONT_H - 0.15) / n
        click_shapes = []
        for i, item in enumerate(word_match):
            y = CONT_Y + 0.08 + i * row_h
            desc = item.get("description", "")
            word = item.get("word", "")
            rect(s, 0.35, y + 0.04, 6.8, row_h - 0.08, C["WHITE"], "CCCCCC", 0.6)
            txt(s, desc, 0.5, y + 0.06, 6.6, row_h - 0.12,
                size=13, color=C["BLACK"], align="left", valign="middle")
            ans_shape = txt(s, word, 7.3, y + 0.10, 2.4, row_h - 0.20,
                size=18, bold=True, color=C["PINK"], align="center", valign="middle")
            click_shapes.append(ans_shape)

        register_clicks(s, click_shapes)

    def slide_y3_spell_check_sentences_blank():
        """Y3 Spell Check and Sentences combined (blank)."""
        spell_data = lesson.get("spellData", [])

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Independent", "Spell Check and Sentences",
                  "Put a ✓ if spelled correctly and a ✗ if not. Correct the mistakes. Then write sentences.",
                  f"{CODE}.15")

        n_words = min(len(spell_data), 10)
        top_h = CONT_H * 0.62
        row_h = top_h / max(n_words, 1)

        hdr = [
            {"text": "Word",            "bold": True, "fill": C["WHITE"], "align": "center", "size": 12},
            {"text": "✓ or ✗",         "bold": True, "fill": C["WHITE"], "align": "center", "size": 12},
            {"text": "Correct spelling","bold": True, "fill": C["WHITE"], "align": "center", "size": 12},
        ]
        rows = []
        for i, row in enumerate(spell_data[:n_words]):
            opts = row.get("opts", [])
            correct_idx = row.get("correct", 0)
            correct_word = opts[correct_idx] if correct_idx < len(opts) else ""
            # Randomly show either correct or one of the incorrect spellings
            wrong_opts = [o for j, o in enumerate(opts) if j != correct_idx]
            display_word = correct_word if i % 2 == 0 else (wrong_opts[0] if wrong_opts else correct_word)
            rows.append([
                {"text": display_word, "fill": C["WHITE"], "align": "center", "size": 13},
                {"text": "",           "fill": C["WHITE"], "align": "center"},
                {"text": "",           "fill": C["WHITE"]},
            ])

        rh_list = [min(0.38, row_h)] * (n_words + 1)
        table(s, [hdr] + rows,
              0.35, CONT_Y + 0.05, 9.3, top_h,
              [4.0, 1.5, 3.8], rh_list)

        # Sentence prompts below
        sent_y = CONT_Y + top_h + 0.10
        sent_h = CONT_H - top_h - 0.15
        txt(s, f"Write the word '{WORDS[0]}' in a sentence.",
            0.5, sent_y, 9.0, sent_h * 0.45,
            size=13, color=C["BLACK"], align="left", valign="top")
        txt(s, "Choose one or more other words to write in a sentence.",
            0.5, sent_y + sent_h * 0.48, 9.0, sent_h * 0.45,
            size=13, color=C["BLACK"], align="left", valign="top")

    def slide_y3_spell_check_sentences_answers():
        """Y3 Spell Check and Sentences answers."""
        spell_data = lesson.get("spellData", [])

        s = prs.slides.add_slide(BLANK)
        add_frame(s, "Whole Group", "Spell Check and Sentences", "Answers", f"{CODE}.16")

        n_words = min(len(spell_data), 10)
        top_h = CONT_H * 0.62
        row_h = top_h / max(n_words, 1)

        hdr = [
            {"text": "Word",            "bold": True, "fill": C["WHITE"], "align": "center", "size": 12},
            {"text": "✓ or ✗",         "bold": True, "fill": C["WHITE"], "align": "center", "size": 12},
            {"text": "Correct spelling","bold": True, "fill": C["WHITE"], "align": "center", "size": 12},
        ]
        rows = []
        for i, row in enumerate(spell_data[:n_words]):
            opts = row.get("opts", [])
            correct_idx = row.get("correct", 0)
            correct_word = opts[correct_idx] if correct_idx < len(opts) else ""
            wrong_opts = [o for j, o in enumerate(opts) if j != correct_idx]
            display_word = correct_word if i % 2 == 0 else (wrong_opts[0] if wrong_opts else correct_word)
            is_shown_correct = (display_word == correct_word)
            rows.append([
                {"text": display_word,  "fill": C["WHITE"], "align": "center", "size": 13},
                {"text": "✓" if is_shown_correct else "✗",
                 "fill": C["WHITE"], "align": "center", "size": 16,
                 "bold": True,
                 "color": C["GREEN_S"] if is_shown_correct else C["RED"]},
                {"text": "" if is_shown_correct else correct_word,
                 "fill": C["WHITE"], "align": "center", "size": 13,
                 "bold": True, "color": C["PINK"]},
            ])

        rh_list = [min(0.38, row_h)] * (n_words + 1)
        table(s, [hdr] + rows,
              0.35, CONT_Y + 0.05, 9.3, top_h,
              [4.0, 1.5, 3.8], rh_list)

        # Sample sentence answer
        sent_y = CONT_Y + top_h + 0.10
        sent_h = CONT_H - top_h - 0.15
        sample = SENTENCES.get(WORDS[0], "")
        txt(s, f"e.g. '{sample}'",
            0.5, sent_y, 9.0, sent_h,
            size=13, italic=True, color=C["PINK"], align="left", valign="middle")

    # ════════════════════════════════════════════════════════════════════════
    # BUILD SEQUENCE — branches on year group
    # ════════════════════════════════════════════════════════════════════════

    if YEAR_GROUP == "Y2":
        slide1()
        slide_key_spelling()
        slide_y2_starter()
        slide4()                          # This Week's Words (shared)
        slide_y2_spelling_pattern()
        if lesson.get("y2IncludeMorphMatrix"):
            slide_y2_morph_matrix_blank()
            slide_y2_morph_matrix_answers()
        slide_y2_word_sort_blank()
        slide_y2_word_sort_answers()
        slide_y2_sound_it_squash_it()
        slide_y2_phoneme_maps_blank()
        slide_y2_phoneme_maps_answers()
        slide_y2_sentences_synonyms_blank()
        slide_y2_words_in_action_blank()
        slide_y2_sentences_synonyms_answers()
        slide_y2_words_in_action_answers()
        slide_y2_word_spotter()
        slide_y2_continuous_provision()

    elif YEAR_GROUP == "Y3":
        slide1()
        slide_key_spelling()
        slide_y3_starter_blank()
        slide_y3_starter_answers()
        slide_y3_this_weeks_words()
        slide_y3_spelling_pattern()
        slide_y3_etymology_simple()
        slide_y3_morph_matrix_blank()
        slide_y3_morph_matrix_answers()
        slide_y3_word_sort_blank()
        slide_y3_word_sort_answers()
        slide_y3_sound_buttons_syllables()
        slide_y3_word_maps_blank()
        slide_y3_word_maps_answers()
        slide_y3_word_match_blank()
        slide_y3_word_match_answers()
        slide_y3_spell_check_sentences_blank()
        slide_y3_spell_check_sentences_answers()

    else:
        # Y4 / Y5 / Y6 — original sequence
        slide1();  slide_key_spelling();  slide2();  slide3();  slide4();  slide5();  slide6();  slide7()
        slide8();  slide9();  slide10(); slide11(); slide12(); slide13(); slide14()
        slide15(); slide16(); slide17(); slide18(); slide19(); slide20(); slide21()

    inject_click_animations()

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    pptx_bytes = buf.read()

    if key_spelling_slide_index[0] is not None and key_spelling_pic_id[0] is not None:
        pptx_bytes = _embed_timer_video(
            pptx_bytes,
            slide_index_zero_based=key_spelling_slide_index[0],
            pic_shape_id=key_spelling_pic_id[0],
        )

    return pptx_bytes


def _embed_timer_video(pptx_bytes: bytes, slide_index_zero_based: int,
                      pic_shape_id: int) -> bytes:
    """Inject the 1-minute timer MP4 into a generated PPTX."""
    import zipfile, io, os
    from lxml import etree

    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
    CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

    mp4_path = os.path.join(os.path.dirname(__file__), "assets", "timer_1min.mp4")
    if not os.path.exists(mp4_path):
        return pptx_bytes
    with open(mp4_path, "rb") as f:
        mp4_bytes = f.read()

    slide_one_based = slide_index_zero_based + 1
    slide_xml_path = f"ppt/slides/slide{slide_one_based}.xml"
    slide_rels_path = f"ppt/slides/_rels/slide{slide_one_based}.xml.rels"

    src_zip = zipfile.ZipFile(io.BytesIO(pptx_bytes), "r")
    namelist = src_zip.namelist()

    ct_xml = src_zip.read("[Content_Types].xml")
    ct_tree = etree.fromstring(ct_xml)
    has_mp4 = False
    for default in ct_tree.findall(f"{{{CT_NS}}}Default"):
        if default.get("Extension") == "mp4":
            has_mp4 = True
            break
    if not has_mp4:
        new_default = etree.SubElement(ct_tree, f"{{{CT_NS}}}Default")
        new_default.set("Extension", "mp4")
        new_default.set("ContentType", "video/mp4")
    new_ct_xml = etree.tostring(ct_tree, xml_declaration=True,
                                encoding="UTF-8", standalone=True)

    rels_xml = src_zip.read(slide_rels_path)
    rels_tree = etree.fromstring(rels_xml)
    existing_ids = set()
    for r in rels_tree.findall(f"{{{REL_NS}}}Relationship"):
        existing_ids.add(r.get("Id"))

    def next_rid():
        i = 1
        while f"rId{i}" in existing_ids:
            i += 1
        existing_ids.add(f"rId{i}")
        return f"rId{i}"

    media_rid = next_rid()
    video_rid = next_rid()

    media_rel = etree.SubElement(rels_tree, f"{{{REL_NS}}}Relationship")
    media_rel.set("Id", media_rid)
    media_rel.set("Type", "http://schemas.microsoft.com/office/2007/relationships/media")
    media_rel.set("Target", "../media/timer_1min.mp4")

    video_rel = etree.SubElement(rels_tree, f"{{{REL_NS}}}Relationship")
    video_rel.set("Id", video_rid)
    video_rel.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video")
    video_rel.set("Target", "../media/timer_1min.mp4")

    new_rels_xml = etree.tostring(rels_tree, xml_declaration=True,
                                  encoding="UTF-8", standalone=True)

    slide_xml = src_zip.read(slide_xml_path)
    slide_tree = etree.fromstring(slide_xml)

    pic = None
    for p in slide_tree.iter(f"{{{P}}}pic"):
        cNvPr = p.find(f".//{{{P}}}cNvPr")
        if cNvPr is not None and cNvPr.get("id") == str(pic_shape_id):
            pic = p
            break
    if pic is None:
        src_zip.close()
        return pptx_bytes

    cNvPr = pic.find(f".//{{{P}}}cNvPr")
    if cNvPr is not None:
        if cNvPr.find(f"{{{A}}}hlinkClick") is None:
            hl = etree.SubElement(cNvPr, f"{{{A}}}hlinkClick")
            hl.set(f"{{{R}}}id", "")
            hl.set("action", "ppaction://media")
            ext = cNvPr.find(f"{{{A}}}extLst")
            if ext is not None:
                cNvPr.remove(hl)
                cNvPr.insert(list(cNvPr).index(ext), hl)

    nvPr = pic.find(f".//{{{P}}}nvPr")
    if nvPr is not None:
        videoFile = etree.SubElement(nvPr, f"{{{A}}}videoFile")
        videoFile.set(f"{{{R}}}link", video_rid)
        ext_list = etree.SubElement(nvPr, f"{{{P}}}extLst")
        ext = etree.SubElement(ext_list, f"{{{P}}}ext")
        ext.set("uri", "{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}")
        media = etree.SubElement(ext, f"{{{P14}}}media")
        media.set(f"{{{R}}}embed", media_rid)

    for tag in ("timing", "bldLst"):
        existing = slide_tree.find(f"{{{P}}}{tag}")
        if existing is not None:
            slide_tree.remove(existing)

    timing = etree.SubElement(slide_tree, f"{{{P}}}timing")
    tnLst = etree.SubElement(timing, f"{{{P}}}tnLst")
    root_par = etree.SubElement(tnLst, f"{{{P}}}par")
    root_cTn = etree.SubElement(root_par, f"{{{P}}}cTn")
    root_cTn.set("id", "1"); root_cTn.set("dur", "indefinite")
    root_cTn.set("restart", "never"); root_cTn.set("nodeType", "tmRoot")
    root_children = etree.SubElement(root_cTn, f"{{{P}}}childTnLst")

    seq = etree.SubElement(root_children, f"{{{P}}}seq")
    seq.set("concurrent", "1"); seq.set("nextAc", "seek")
    seq_cTn = etree.SubElement(seq, f"{{{P}}}cTn")
    seq_cTn.set("id", "2"); seq_cTn.set("dur", "indefinite"); seq_cTn.set("nodeType", "mainSeq")
    seq_children = etree.SubElement(seq_cTn, f"{{{P}}}childTnLst")

    click_par = etree.SubElement(seq_children, f"{{{P}}}par")
    click_cTn = etree.SubElement(click_par, f"{{{P}}}cTn")
    click_cTn.set("id", "3"); click_cTn.set("fill", "hold")
    cl_stCondLst = etree.SubElement(click_cTn, f"{{{P}}}stCondLst")
    etree.SubElement(cl_stCondLst, f"{{{P}}}cond").set("delay", "indefinite")
    cl_children = etree.SubElement(click_cTn, f"{{{P}}}childTnLst")

    inner_par = etree.SubElement(cl_children, f"{{{P}}}par")
    inner_cTn = etree.SubElement(inner_par, f"{{{P}}}cTn")
    inner_cTn.set("id", "4"); inner_cTn.set("fill", "hold")
    in_stCondLst = etree.SubElement(inner_cTn, f"{{{P}}}stCondLst")
    etree.SubElement(in_stCondLst, f"{{{P}}}cond").set("delay", "0")
    in_children = etree.SubElement(inner_cTn, f"{{{P}}}childTnLst")

    play_par = etree.SubElement(in_children, f"{{{P}}}par")
    play_cTn = etree.SubElement(play_par, f"{{{P}}}cTn")
    play_cTn.set("id", "5"); play_cTn.set("presetID", "1")
    play_cTn.set("presetClass", "mediacall"); play_cTn.set("presetSubtype", "0")
    play_cTn.set("fill", "hold"); play_cTn.set("nodeType", "clickEffect")
    play_stCondLst = etree.SubElement(play_cTn, f"{{{P}}}stCondLst")
    etree.SubElement(play_stCondLst, f"{{{P}}}cond").set("delay", "0")
    play_children = etree.SubElement(play_cTn, f"{{{P}}}childTnLst")

    cmd = etree.SubElement(play_children, f"{{{P}}}cmd")
    cmd.set("type", "call"); cmd.set("cmd", "playFrom(0.0)")
    cBhvr = etree.SubElement(cmd, f"{{{P}}}cBhvr")
    cBhvr_cTn = etree.SubElement(cBhvr, f"{{{P}}}cTn")
    cBhvr_cTn.set("id", "6"); cBhvr_cTn.set("dur", "63000"); cBhvr_cTn.set("fill", "hold")
    tgtEl = etree.SubElement(cBhvr, f"{{{P}}}tgtEl")
    spTgt = etree.SubElement(tgtEl, f"{{{P}}}spTgt")
    spTgt.set("spid", str(pic_shape_id))

    prev = etree.SubElement(seq, f"{{{P}}}prevCondLst")
    pc = etree.SubElement(prev, f"{{{P}}}cond")
    pc.set("evt", "onPrev"); pc.set("delay", "0")
    pc_tgt = etree.SubElement(pc, f"{{{P}}}tgtEl")
    etree.SubElement(pc_tgt, f"{{{P}}}sldTgt")
    nxt = etree.SubElement(seq, f"{{{P}}}nextCondLst")
    nc = etree.SubElement(nxt, f"{{{P}}}cond")
    nc.set("evt", "onNext"); nc.set("delay", "0")
    nc_tgt = etree.SubElement(nc, f"{{{P}}}tgtEl")
    etree.SubElement(nc_tgt, f"{{{P}}}sldTgt")

    video_node = etree.SubElement(root_children, f"{{{P}}}video")
    cMediaNode = etree.SubElement(video_node, f"{{{P}}}cMediaNode")
    cMediaNode.set("vol", "80000")
    media_cTn = etree.SubElement(cMediaNode, f"{{{P}}}cTn")
    media_cTn.set("id", "7"); media_cTn.set("fill", "hold"); media_cTn.set("display", "0")
    m_stCondLst = etree.SubElement(media_cTn, f"{{{P}}}stCondLst")
    etree.SubElement(m_stCondLst, f"{{{P}}}cond").set("delay", "indefinite")
    m_tgtEl = etree.SubElement(cMediaNode, f"{{{P}}}tgtEl")
    m_spTgt = etree.SubElement(m_tgtEl, f"{{{P}}}spTgt")
    m_spTgt.set("spid", str(pic_shape_id))

    iseq = etree.SubElement(root_children, f"{{{P}}}seq")
    iseq.set("concurrent", "1"); iseq.set("nextAc", "seek")
    iseq_cTn = etree.SubElement(iseq, f"{{{P}}}cTn")
    iseq_cTn.set("id", "8"); iseq_cTn.set("restart", "whenNotActive")
    iseq_cTn.set("fill", "hold"); iseq_cTn.set("evtFilter", "cancelBubble")
    iseq_cTn.set("nodeType", "interactiveSeq")
    iseq_stCondLst = etree.SubElement(iseq_cTn, f"{{{P}}}stCondLst")
    iseq_cond = etree.SubElement(iseq_stCondLst, f"{{{P}}}cond")
    iseq_cond.set("evt", "onClick"); iseq_cond.set("delay", "0")
    iseq_cond_tgt = etree.SubElement(iseq_cond, f"{{{P}}}tgtEl")
    iseq_cond_sp = etree.SubElement(iseq_cond_tgt, f"{{{P}}}spTgt")
    iseq_cond_sp.set("spid", str(pic_shape_id))
    iseq_endSync = etree.SubElement(iseq_cTn, f"{{{P}}}endSync")
    iseq_endSync.set("evt", "end"); iseq_endSync.set("delay", "0")
    iseq_rtn = etree.SubElement(iseq_endSync, f"{{{P}}}rtn")
    iseq_rtn.set("val", "all")
    iseq_children = etree.SubElement(iseq_cTn, f"{{{P}}}childTnLst")

    toggle_par = etree.SubElement(iseq_children, f"{{{P}}}par")
    toggle_cTn = etree.SubElement(toggle_par, f"{{{P}}}cTn")
    toggle_cTn.set("id", "9"); toggle_cTn.set("fill", "hold")
    t_stCondLst = etree.SubElement(toggle_cTn, f"{{{P}}}stCondLst")
    etree.SubElement(t_stCondLst, f"{{{P}}}cond").set("delay", "0")
    t_children = etree.SubElement(toggle_cTn, f"{{{P}}}childTnLst")

    t_inner_par = etree.SubElement(t_children, f"{{{P}}}par")
    t_inner_cTn = etree.SubElement(t_inner_par, f"{{{P}}}cTn")
    t_inner_cTn.set("id", "10"); t_inner_cTn.set("fill", "hold")
    ti_stCondLst = etree.SubElement(t_inner_cTn, f"{{{P}}}stCondLst")
    etree.SubElement(ti_stCondLst, f"{{{P}}}cond").set("delay", "0")
    ti_children = etree.SubElement(t_inner_cTn, f"{{{P}}}childTnLst")

    pause_par = etree.SubElement(ti_children, f"{{{P}}}par")
    pause_cTn = etree.SubElement(pause_par, f"{{{P}}}cTn")
    pause_cTn.set("id", "11"); pause_cTn.set("presetID", "2")
    pause_cTn.set("presetClass", "mediacall"); pause_cTn.set("presetSubtype", "0")
    pause_cTn.set("fill", "hold"); pause_cTn.set("nodeType", "clickEffect")
    pause_stCondLst = etree.SubElement(pause_cTn, f"{{{P}}}stCondLst")
    etree.SubElement(pause_stCondLst, f"{{{P}}}cond").set("delay", "0")
    pause_children = etree.SubElement(pause_cTn, f"{{{P}}}childTnLst")

    pause_cmd = etree.SubElement(pause_children, f"{{{P}}}cmd")
    pause_cmd.set("type", "call"); pause_cmd.set("cmd", "togglePause")
    pause_cBhvr = etree.SubElement(pause_cmd, f"{{{P}}}cBhvr")
    pause_cBhvr_cTn = etree.SubElement(pause_cBhvr, f"{{{P}}}cTn")
    pause_cBhvr_cTn.set("id", "12"); pause_cBhvr_cTn.set("dur", "1"); pause_cBhvr_cTn.set("fill", "hold")
    pause_tgtEl = etree.SubElement(pause_cBhvr, f"{{{P}}}tgtEl")
    pause_spTgt = etree.SubElement(pause_tgtEl, f"{{{P}}}spTgt")
    pause_spTgt.set("spid", str(pic_shape_id))

    iseq_next = etree.SubElement(iseq, f"{{{P}}}nextCondLst")
    iseq_next_cond = etree.SubElement(iseq_next, f"{{{P}}}cond")
    iseq_next_cond.set("evt", "onClick"); iseq_next_cond.set("delay", "0")
    iseq_next_tgt = etree.SubElement(iseq_next_cond, f"{{{P}}}tgtEl")
    iseq_next_sp = etree.SubElement(iseq_next_tgt, f"{{{P}}}spTgt")
    iseq_next_sp.set("spid", str(pic_shape_id))

    new_slide_xml = etree.tostring(slide_tree, xml_declaration=True,
                                   encoding="UTF-8", standalone=True)

    out_buf = io.BytesIO()
    out_zip = zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED)
    for name in namelist:
        if name == "[Content_Types].xml":
            out_zip.writestr(name, new_ct_xml)
        elif name == slide_xml_path:
            out_zip.writestr(name, new_slide_xml)
        elif name == slide_rels_path:
            out_zip.writestr(name, new_rels_xml)
        else:
            out_zip.writestr(name, src_zip.read(name))
    if "ppt/media/timer_1min.mp4" not in namelist:
        out_zip.writestr("ppt/media/timer_1min.mp4", mp4_bytes)
    out_zip.close()
    src_zip.close()
    return out_buf.getvalue()
